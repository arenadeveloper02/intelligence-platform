"""Platform hub server — Google Sign-In + multi-dashboard routing."""

import os
import json
import uuid
import logging
import re
from datetime import datetime, timezone, timedelta
from pathlib import Path
from functools import wraps

from flask import (
    Flask, send_file, send_from_directory, abort, jsonify,
    request, session, redirect, url_for,
    make_response, render_template,
)
import requests
from collections import Counter

log = logging.getLogger(__name__)

# Indian Standard Time = UTC+5:30
IST = timezone(timedelta(hours=5, minutes=30))

app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "cst-dev-secret-do-not-use-in-prod-abc123xyz")
app.permanent_session_lifetime = timedelta(days=7)

@app.errorhandler(403)
def forbidden(e):
    return render_template("403.html"), 403

@app.errorhandler(500)
def server_error(e):
    """Always answer API routes with JSON so the frontend never chokes on an HTML error page."""
    if request.path.startswith("/api/"):
        return jsonify({"error": "Kairo is taking longer than usual - please try again."}), 500
    return ("Internal Server Error", 500)

# ── Google OAuth ────────────────────────────────────────────────────────────────
# Set GOOGLE_CLIENT_ID in Railway → Variables.
# Setup: console.cloud.google.com → APIs & Services → Credentials
#        → Create OAuth 2.0 Client ID → Web application
#        → Authorised JavaScript origins: https://intelligence.position2.com
GOOGLE_CLIENT_ID = os.environ.get("GOOGLE_CLIENT_ID", "")

# Anonymous Visitors Google Sheet
ANON_VISITORS_SHEET_ID = "1y5_ef9Df5v8PVuGs60DzzlzE1ZJLxwvB5MQsB6ug458"

# Google Sheet ID for login tracking (set LOGIN_LOG_SHEET_ID in Railway Variables).
# Create a new Google Sheet, share it with the service account email (Editor),
# then paste the sheet ID here (from its URL: /spreadsheets/d/<ID>/edit).
LOGIN_LOG_SHEET_ID = os.environ.get("LOGIN_LOG_SHEET_ID", "")
_SA_JSON = str(Path(__file__).parent / "service_account.json")

# ── Login logger ────────────────────────────────────────────────────────────────
def _parse_ua(ua: str) -> tuple[str, str, str, str]:
    """Return (browser_name, browser_version, os_name, device_type) from User-Agent."""
    ua = ua or ""
    # Device type
    if re.search(r"Mobile|Android|iPhone|iPod", ua, re.I):
        device = "Mobile"
    elif re.search(r"iPad|Tablet", ua, re.I):
        device = "Tablet"
    else:
        device = "Desktop"
    # OS
    if re.search(r"Windows NT", ua):
        os_name = "Windows"
    elif re.search(r"Mac OS X", ua):
        os_name = "macOS"
    elif re.search(r"Android", ua):
        os_name = "Android"
    elif re.search(r"iPhone|iPad", ua):
        os_name = "iOS"
    elif re.search(r"Linux", ua):
        os_name = "Linux"
    else:
        os_name = "Unknown"
    # Browser (order matters — Chrome must come before Safari)
    m = re.search(r"Edg(?:e)?/([\d.]+)", ua)
    if m: return "Edge", m.group(1), os_name, device
    m = re.search(r"OPR/([\d.]+)", ua)
    if m: return "Opera", m.group(1), os_name, device
    m = re.search(r"Firefox/([\d.]+)", ua)
    if m: return "Firefox", m.group(1), os_name, device
    m = re.search(r"Chrome/([\d.]+)", ua)
    if m: return "Chrome", m.group(1), os_name, device
    m = re.search(r"Version/([\d.]+).*Safari", ua)
    if m: return "Safari", m.group(1), os_name, device
    return "Unknown", "", os_name, device


def _log_login_to_sheet(user: dict) -> None:
    """Append one login row to the tracking Google Sheet. Fails silently."""
    if not LOGIN_LOG_SHEET_ID:
        return
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
        import json as _json

        # Prefer env var (Railway) over local file
        sa_json_str = os.environ.get("GOOGLE_SA_JSON", "")
        if sa_json_str:
            sa_info = _json.loads(sa_json_str)
            creds = service_account.Credentials.from_service_account_info(
                sa_info,
                scopes=["https://www.googleapis.com/auth/spreadsheets"],
            )
        elif Path(_SA_JSON).exists():
            creds = service_account.Credentials.from_service_account_file(
                _SA_JSON,
                scopes=["https://www.googleapis.com/auth/spreadsheets"],
            )
        else:
            log.warning("Login sheet: no credentials found (set GOOGLE_SA_JSON env var)")
            return
        svc = build("sheets", "v4", credentials=creds, cache_discovery=False)

        now = datetime.now(IST)
        ua_raw  = request.headers.get("User-Agent", "")
        browser, bv, os_name, device = _parse_ua(ua_raw)
        ip = (request.headers.get("X-Forwarded-For", "") or
              request.headers.get("X-Real-IP", "") or
              request.remote_addr or "")
        ip = ip.split(",")[0].strip()  # X-Forwarded-For can be a list

        # 20 columns — add header row automatically on first write
        row = [
            now.strftime("%Y-%m-%d %H:%M:%S IST"),   # 1  Timestamp
            now.strftime("%Y-%m-%d"),                  # 2  Date
            now.strftime("%H:%M:%S"),                  # 3  Time (IST)
            now.strftime("%A"),                         # 4  Day of Week
            now.strftime("%H"),                         # 5  Hour (0-23, IST)
            user.get("email", ""),                      # 6  Email
            user.get("name", ""),                       # 7  Full Name
            user.get("given_name", ""),                 # 8  First Name
            user.get("picture", ""),                    # 9  Profile Picture URL
            ip,                                         # 10 IP Address
            browser,                                    # 11 Browser
            bv,                                         # 12 Browser Version
            os_name,                                    # 13 Operating System
            device,                                     # 14 Device Type
            ua_raw[:200],                               # 15 User Agent (truncated)
            request.referrer or "direct",               # 16 Referrer
            "/hub",                                     # 17 Landing Page
            "Google OAuth",                             # 18 Auth Method
            str(uuid.uuid4())[:8],                      # 19 Session ID (short)
            "intelligence.position2.com",               # 20 Platform
        ]

        # Check if header row exists; if sheet is empty, prepend it
        result = svc.spreadsheets().values().get(
            spreadsheetId=LOGIN_LOG_SHEET_ID, range="A1:A1"
        ).execute()
        if not result.get("values"):
            header = [[
                "Timestamp (IST)", "Date", "Time (IST)", "Day of Week", "Hour (IST)",
                "Email", "Full Name", "First Name", "Profile Picture",
                "IP Address", "Browser", "Browser Version", "OS", "Device",
                "User Agent", "Referrer", "Landing Page", "Auth Method",
                "Session ID", "Platform",
            ]]
            svc.spreadsheets().values().append(
                spreadsheetId=LOGIN_LOG_SHEET_ID,
                range="A1",
                valueInputOption="RAW",
                body={"values": header},
            ).execute()

        svc.spreadsheets().values().append(
            spreadsheetId=LOGIN_LOG_SHEET_ID,
            range="A1",
            valueInputOption="RAW",
            insertDataOption="INSERT_ROWS",
            body={"values": [row]},
        ).execute()

    except Exception as e:
        log.warning("Login sheet log failed: %s", e)


# ── Account registry ────────────────────────────────────────────────────────────
ACCOUNTS = {
    "healthcare": {
        "name":        "Healthcare",
        "description": "1,251 healthcare companies tracked for funding, C-suite moves, M&A, and news signals.",
        "icon":        "🏥",
        "accent":      "#3b82f6",
        "dashboard":   Path(__file__).parent / "reports" / "dashboard.html",
    },
    "csg": {
        "name":        "CSG",
        "description": "CSG company intelligence — funding rounds, leadership changes, and market signals.",
        "icon":        "📡",
        "accent":      "#8b5cf6",
        "dashboard":   Path(__file__).parent / "reports" / "dashboard_csg.html",
    },
}

# ── Auth helpers ────────────────────────────────────────────────────────────────
ADMIN_EMAILS = {"krishna.ladha@position2.com", "sudheer.d@position2.com"}

def _get_user():
    """Return current user dict or None."""
    return session.get("google_user")

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not _get_user():
            return redirect(url_for("login_page") + "?error=Your+session+expired.+Please+sign+in+again.")
        return f(*args, **kwargs)
    return decorated

def admin_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        user = _get_user()
        if not user:
            return redirect(url_for("login_page") + "?error=Your+session+expired.+Please+sign+in+again.")
        if user.get("email", "").lower() not in ADMIN_EMAILS:
            abort(403)
        return f(*args, **kwargs)
    return decorated

# ── Google Sign-In ──────────────────────────────────────────────────────────────
@app.route("/auth/google", methods=["POST"])
def auth_google():
    credential = (request.json or {}).get("credential", "")
    if not credential:
        return jsonify({"success": False, "error": "No credential"}), 400

    if not GOOGLE_CLIENT_ID:
        # Dev mode: decode without verification (localhost only)
        import base64, json as _j
        try:
            pad = credential.split(".")[1]
            pad += "=" * (-len(pad) % 4)
            idinfo = _j.loads(base64.urlsafe_b64decode(pad))
        except Exception as e:
            return jsonify({"success": False, "error": str(e)}), 401
    else:
        try:
            from google.oauth2 import id_token
            from google.auth.transport import requests as greq
            idinfo = id_token.verify_oauth2_token(credential, greq.Request(), GOOGLE_CLIENT_ID)
        except Exception as e:
            return jsonify({"success": False, "error": str(e)}), 401

    email = idinfo.get("email", "")
    if not email.lower().endswith("@position2.com"):
        return jsonify({"success": False, "error": "Access restricted to Position2 accounts only."}), 403

    session["google_user"] = {
        "email":      email,
        "name":       idinfo.get("name", ""),
        "given_name": idinfo.get("given_name", ""),
        "picture":    idinfo.get("picture", ""),
    }
    session.permanent = True
    _log_login_to_sheet(session["google_user"])   # fire-and-forget, fails silently
    return jsonify({"success": True, "redirect": "/hub"})


# ── Core routes ─────────────────────────────────────────────────────────────────
@app.route("/")
def index():
    return redirect(url_for("hub") if _get_user() else url_for("login_page"))

@app.route("/login")
def login_page():
    if _get_user():
        return redirect(url_for("hub"))
    return render_template("login.html", google_client_id=GOOGLE_CLIENT_ID,
                           error=request.args.get("error", ""))

@app.route("/login-preview")
def login_preview():
    return render_template("login_preview.html", google_client_id=GOOGLE_CLIENT_ID,
                           error=request.args.get("error", ""))

@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login_page"))

# ── Hub pages ───────────────────────────────────────────────────────────────────
@app.route("/hub")
@login_required
def hub():
    return render_template("hub.html", user=_get_user())

@app.route("/ppc")
@login_required
def ppc():
    return render_template("ppc.html", user=_get_user())

@app.route("/seo")
@login_required
def seo():
    return render_template("seo.html", user=_get_user())

# ── Embedded dashboards ─────────────────────────────────────────────────────────
_SERP_BASE = "https://serp-content-researcher-production-a947.up.railway.app"

# ── Ad Intelligence (built React app served directly — no iframe) ────────────
AD_INTEL_SHEET_ID = "16U5_QSxMmrAGKvK5dHScBu1Et4BJ1p8Q1ns5LycRA0s"

@app.route("/ppc/ad-intelligence")
@app.route("/ppc/ad-intelligence/")
@login_required
def ad_intelligence():
    return send_from_directory("ad_intelligence", "index.html")

@app.route("/ppc/ad-intelligence/assets/<path:filename>")
def ad_intelligence_assets(filename):
    return send_from_directory("ad_intelligence/assets", filename)

@app.route("/ppc/ad-intelligence/favicon.svg")
def ad_intelligence_favicon():
    return send_from_directory("ad_intelligence", "favicon.svg")

@app.route("/ppc/ad-intelligence/icons.svg")
def ad_intelligence_icons():
    return send_from_directory("ad_intelligence", "icons.svg")

_SEO_TOOLS = {
    "seo-geo-audit":          ("/seo-geo-audit",          "SEO & GEO Audit"),
    "article-recommendation": ("/article-recommendation", "Article Recommendation"),
    "content-enhancement":    ("/content-enhancement",    "Content Enhancement"),
    "content-research":       ("/content-research",       "Content Research"),
    "keyword-research":       ("/keyword-research",       "Keyword Research"),
    "agent-readiness-audit":  ("/agent-readiness-audit",  "Agent Readiness Audit"),
    "image-alt-audit":        ("/image-alt-audit",        "Image Alt Tag Audit"),
    "knowledge-base":         ("/kb",                     "Knowledge Base"),
    "team-insights":          ("/team-insights",          "Team Insights"),
    "serp-researcher":        ("/",                       "SERP Content Researcher"),
}

@app.route("/seo/<tool_slug>")
@login_required
def seo_tool(tool_slug: str):
    if tool_slug not in _SEO_TOOLS:
        abort(404)
    tool_path, tool_name = _SEO_TOOLS[tool_slug]
    pt = os.environ.get("SERP_PLATFORM_TOKEN", "")
    sep = "?" if "?" not in tool_path else "&"
    embed_url = f"{_SERP_BASE}{tool_path}{sep + 'pt=' + pt if pt else ''}"
    return render_template("embed.html",
        user=_get_user(),
        title=tool_name,
        embed_url=embed_url,
        breadcrumb=[("Hub", "/hub"), ("SEO", "/seo")],
        current=tool_name,
        accent="#34d399",
    )

# ── Company Signal Tracker ───────────────────────────────────────────────────────
@app.route("/accounts")
@login_required
def accounts():
    cards_html = "".join(_build_account_card(aid, cfg) for aid, cfg in ACCOUNTS.items())
    return render_template("accounts.html", user=_get_user(), account_cards=cards_html)

@app.route("/signal-tracker/<account_id>")
@app.route("/signal-tracker/<account_id>/<section>")
@login_required
def dashboard(account_id: str, section: str = None):
    cfg = ACCOUNTS.get(account_id)
    if not cfg:
        abort(404, f"Unknown account '{account_id}'")
    path: Path = cfg["dashboard"]
    if not path.exists():
        abort(404, f"Dashboard for '{cfg['name']}' not generated yet.")
    resp = make_response(send_file(str(path)))
    resp.headers.update({"Cache-Control": "no-cache, no-store, must-revalidate",
                         "Pragma": "no-cache", "Expires": "0"})
    return resp

@app.after_request
def _no_html_cache(resp):
    """Never let browsers cache HTML pages — UI updates must show immediately after deploys."""
    try:
        if resp.mimetype == "text/html":
            resp.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
            resp.headers["Pragma"] = "no-cache"
            resp.headers["Expires"] = "0"
    except Exception:
        pass
    return resp


@app.route("/dashboard/<account_id>")
@app.route("/dashboard/<account_id>/<section>")
@login_required
def dashboard_legacy(account_id: str, section: str = None):
    """Back-compat: old /dashboard/* URLs redirect to canonical /signal-tracker/*."""
    target = "/signal-tracker/" + account_id + (("/" + section) if section else "")
    return redirect(target, code=301)

@app.route("/api/whoami")
@login_required
def whoami():
    u = _get_user() or {}
    return jsonify({"name": u.get("name", ""), "given_name": u.get("given_name", ""),
                    "email": u.get("email", ""), "picture": u.get("picture", "")})

# ── Health + API ─────────────────────────────────────────────────────────────────
@app.route("/api/track", methods=["POST"])
def track_page():
    """Record page view duration to Google Sheet 'Page Views' tab."""
    try:
        # sendBeacon sends text/plain, not application/json — handle both
        data = request.json
        if data is None:
            try:
                data = json.loads(request.get_data(as_text=True))
            except Exception:
                data = {}
        data    = data or {}
        page    = data.get("page", "unknown")
        seconds = int(data.get("seconds", 0))
        email   = data.get("email", "") or (session.get("google_user") or {}).get("email", "")
        title   = data.get("title", page)
        if seconds < 1:
            return jsonify({"ok": True})

        mins, secs = divmod(seconds, 60)
        duration_fmt = f"{mins}m {secs}s" if mins else f"{secs}s"

        now = datetime.now(IST)
        row = [
            now.strftime("%Y-%m-%d %H:%M:%S IST"),
            now.strftime("%Y-%m-%d"),
            now.strftime("%H:%M:%S"),
            now.strftime("%A"),
            email,
            title,
            page,
            seconds,
            duration_fmt,
            (request.headers.get("X-Forwarded-For","") or request.remote_addr or "").split(",")[0].strip(),
            _parse_ua(request.headers.get("User-Agent",""))[0],   # browser
            _parse_ua(request.headers.get("User-Agent",""))[2],   # OS
            _parse_ua(request.headers.get("User-Agent",""))[3],   # device
        ]

        if not LOGIN_LOG_SHEET_ID:
            return jsonify({"ok": True})

        import json as _j
        from google.oauth2 import service_account
        from googleapiclient.discovery import build

        sa_str = os.environ.get("GOOGLE_SA_JSON","")
        if not sa_str:
            return jsonify({"ok": True})

        sa_info = _j.loads(sa_str)
        creds   = service_account.Credentials.from_service_account_info(
            sa_info, scopes=["https://www.googleapis.com/auth/spreadsheets"])
        svc = build("sheets", "v4", credentials=creds, cache_discovery=False)

        # Auto-create header on first write to Page Views tab
        try:
            existing = svc.spreadsheets().values().get(
                spreadsheetId=LOGIN_LOG_SHEET_ID, range="Page Views!A1:A1").execute()
            if not existing.get("values"):
                raise Exception("empty")
        except Exception:
            header = [["Timestamp (IST)","Date","Time (IST)","Day","Email","Page Title",
                       "Page URL","Seconds","Duration","IP","Browser","OS","Device"]]
            try:
                svc.spreadsheets().batchUpdate(
                    spreadsheetId=LOGIN_LOG_SHEET_ID,
                    body={"requests":[{"addSheet":{"properties":{"title":"Page Views"}}}]}
                ).execute()
            except Exception:
                pass
            svc.spreadsheets().values().append(
                spreadsheetId=LOGIN_LOG_SHEET_ID, range="Page Views!A1",
                valueInputOption="RAW", body={"values": header}).execute()

        svc.spreadsheets().values().append(
            spreadsheetId=LOGIN_LOG_SHEET_ID, range="Page Views!A1",
            valueInputOption="RAW", insertDataOption="INSERT_ROWS",
            body={"values": [row]}).execute()

    except Exception as e:
        log.warning("Page track failed: %s", e)

    return jsonify({"ok": True})


def _fetch_usage_data() -> dict:
    """Fetch login + page view data from Sheets. Shared by shell and data endpoints."""
    def _fetch(tab_range):
        try:
            import json as _j
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
            sa_str = os.environ.get("GOOGLE_SA_JSON", "")
            if not sa_str or not LOGIN_LOG_SHEET_ID:
                return []
            sa_info = _j.loads(sa_str)
            creds = service_account.Credentials.from_service_account_info(
                sa_info, scopes=["https://www.googleapis.com/auth/spreadsheets"])
            svc = build("sheets", "v4", credentials=creds, cache_discovery=False)
            r = svc.spreadsheets().values().get(
                spreadsheetId=LOGIN_LOG_SHEET_ID, range=tab_range).execute()
            return r.get("values", [])
        except Exception as e:
            log.warning("admin_usage sheet read failed: %s", e)
            return []

    def col(row, i, default=""):
        return row[i] if len(row) > i else default

    login_rows = _fetch("A1:T1000")
    page_rows  = _fetch("Page Views!A1:M2000")
    login_data = login_rows[1:] if len(login_rows) > 1 else []
    page_data  = page_rows[1:]  if len(page_rows)  > 1 else []

    from collections import Counter

    unique_users     = len({col(r, 5) for r in login_data if col(r, 5)})
    total_logins     = len(login_data)
    total_page_views = len(page_data)

    # Total time spent across all page views
    total_secs = sum(int(col(r, 7)) for r in page_data if col(r, 7).isdigit())
    h, rem = divmod(total_secs, 3600)
    m = rem // 60
    total_time_fmt = f"{h}h {m}m" if h else (f"{m}m" if m else "—")

    # Top pages
    page_counts: dict = {}
    for r in page_data:
        t = col(r, 5)
        if t:
            page_counts[t] = page_counts.get(t, 0) + 1
    top_pages = sorted(page_counts.items(), key=lambda x: x[1], reverse=True)[:15]

    # Logins per day (last 14 days)
    login_days = Counter(col(r, 1) for r in login_data if col(r, 1))
    sorted_days = sorted(login_days.items())[-14:]

    # Browser breakdown (from logins)
    browser_counts = Counter(col(r, 10) for r in login_data if col(r, 10))
    browser_breakdown = browser_counts.most_common(5)

    # Device / OS breakdown + quick facts (from page views)
    device_breakdown = Counter(col(r,12) for r in page_data if col(r,12)).most_common(5)
    os_breakdown     = Counter(col(r,11) for r in page_data if col(r,11)).most_common(5)
    _day_counts      = Counter(col(r,3) for r in page_data if col(r,3))
    busiest_day      = list(_day_counts.most_common(1)[0]) if _day_counts else ["—", 0]
    _avg             = round(total_secs/total_page_views) if total_page_views else 0
    _am, _as         = divmod(_avg, 60)
    avg_view_fmt     = (f"{_am}m {_as}s" if _am else f"{_as}s") if _avg else "—"
    views_per_user   = round(total_page_views/unique_users, 1) if unique_users else 0

    # Per-user activity
    user_map: dict = {}
    for r in login_data:
        e = col(r, 5)
        if not e: continue
        if e not in user_map:
            user_map[e] = {"email": e, "name": col(r, 6), "logins": 0,
                           "last_seen": col(r, 0), "total_secs": 0}
        user_map[e]["logins"] += 1
        user_map[e]["last_seen"] = col(r, 0)   # rows are oldest→newest; last row = most recent
    for r in page_data:
        e = col(r, 4)
        if e in user_map and col(r, 7).isdigit():
            user_map[e]["total_secs"] += int(col(r, 7))
    for u in user_map.values():
        s = u["total_secs"]; uh, ur = divmod(s, 3600); um = ur // 60
        u["time_fmt"] = f"{uh}h {um}m" if uh else (f"{um}m" if um else "—")
    user_activity = sorted(user_map.values(), key=lambda x: x["logins"], reverse=True)

    login_table = [{"ts": col(r,0), "email": col(r,5), "name": col(r,6),
                    "browser": col(r,10), "os": col(r,12), "device": col(r,13)}
                   for r in reversed(login_data)][:100]
    page_table  = [{"ts": col(r,0), "email": col(r,4), "title": col(r,5),
                    "url": col(r,6), "duration": col(r,8)}
                   for r in reversed(page_data)]

    return dict(total_logins=total_logins, unique_users=unique_users,
                total_page_views=total_page_views, total_time_fmt=total_time_fmt,
                top_pages=top_pages, login_days=sorted_days,
                browser_breakdown=browser_breakdown, user_activity=user_activity,
                login_table=login_table, page_table=page_table,
                device_breakdown=device_breakdown, os_breakdown=os_breakdown,
                busiest_day=busiest_day, avg_view_fmt=avg_view_fmt,
                views_per_user=views_per_user)


@app.route("/admin/usage")
@admin_required
def admin_usage():
    """Shell page — renders instantly, JS fetches /admin/usage/data async."""
    return render_template("admin_usage.html", user=_get_user())


@app.route("/admin/usage/data")
@admin_required
def admin_usage_data():
    """JSON data endpoint called by the admin usage shell page."""
    data = _fetch_usage_data()
    return jsonify(data)


def _clean_industry(raw: str) -> str:
    """Strip JSON array brackets from industry field, return first value only."""
    if not raw or raw == "Unavailable":
        return "—"
    if raw.startswith('['):
        try:
            import json as _j
            lst = _j.loads(raw)
            return lst[0].strip() if lst else raw
        except Exception:
            # fallback: strip brackets and quotes manually
            return raw.strip('[]').split(',')[0].strip().strip('"\'')
    return raw


def _fetch_anon_visitors_data() -> dict:
    """Fetch people + company data from the Anonymous Visitors Google Sheet."""
    def _fetch(tab_range):
        try:
            import json as _j
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
            sa_str = os.environ.get("GOOGLE_SA_JSON", "")
            if not sa_str:
                return []
            sa_info = _j.loads(sa_str)
            creds = service_account.Credentials.from_service_account_info(
                sa_info, scopes=["https://www.googleapis.com/auth/spreadsheets.readonly"])
            svc = build("sheets", "v4", credentials=creds, cache_discovery=False)
            r = svc.spreadsheets().values().get(
                spreadsheetId=ANON_VISITORS_SHEET_ID, range=tab_range).execute()
            return r.get("values", [])
        except Exception as e:
            log.warning("anon_visitors sheet read failed: %s", e)
            return []

    def col(row, i, default=""):
        return row[i] if len(row) > i else default

    people_rows  = _fetch("People Enriched!A:K")
    company_rows = _fetch("Visitors By Company!A:J")

    people_data  = people_rows[1:]  if len(people_rows)  > 1 else []
    company_data = company_rows[1:] if len(company_rows) > 1 else []

    from collections import Counter
    ind_counter = Counter()
    for r in company_data:
        ind = col(r, 7)
        if ind and ind != "Unavailable":
            ind_counter[ind] += 1
    top_industries = ind_counter.most_common(8)

    # Deduplicated companies list
    seen, company_table = set(), []
    for r in company_data:
        name = col(r, 0)
        if not name or name in seen:
            continue
        seen.add(name)
        company_table.append({
            "name":      name,
            "website":   col(r, 2),
            "city":      col(r, 3),
            "state":     col(r, 4),
            "country":   col(r, 5),
            "industry":  col(r, 7),
            "employees": col(r, 8),
            "revenue":   col(r, 9),
        })
    company_table.sort(key=lambda x: x["name"])

    # People list — sorted newest first
    people_table = []
    for r in people_data:
        name = col(r, 0)
        if not name or name == "Unavailable":
            continue
        time_str = col(r, 6)
        people_table.append({
            "name":     name,
            "title":    col(r, 1),
            "email":    col(r, 2),
            "location": col(r, 4),
            "pages":    col(r, 5),
            "industry": _clean_industry(col(r, 8)),
            "website":  col(r, 10),
            "date":     time_str[:10] if time_str else "",
            "time_raw": time_str,
        })
    people_table.sort(key=lambda x: x.get("time_raw", ""), reverse=True)

    return dict(
        total_people=len(people_table),
        unique_companies=len(company_table),
        top_industries=top_industries,
        people_table=people_table,
        company_table=company_table,
    )


@app.route("/ppc/anonymous-visitors")
@login_required
def anonymous_visitors():
    """Anonymous Visitors dashboard shell — loads data async."""
    return render_template("anonymous_visitors.html", user=_get_user())


@app.route("/ppc/anonymous-visitors/data")
@login_required
def anonymous_visitors_data():
    """JSON data endpoint for the Anonymous Visitors dashboard."""
    return jsonify(_fetch_anon_visitors_data())




@app.route("/ppc/linkedin-scraper")
@login_required
def linkedin_scraper():
    """LinkedIn ABM Intelligence dashboard — Post & People Intelligence."""
    return render_template("linkedin_scraper.html", user=_get_user())

@app.route("/health")
def health():
    return jsonify({"status": "ok", "accounts": {
        aid: {"name": cfg["name"], "dashboard_exists": cfg["dashboard"].exists()}
        for aid, cfg in ACCOUNTS.items()
    }})

@app.route("/api/weekly-stats")
@app.route("/api/weekly-stats/<account_id>")
@login_required
def weekly_stats(account_id: str = "healthcare"):
    cfg = ACCOUNTS.get(account_id)
    if not cfg:
        return jsonify({"error": f"Unknown account '{account_id}'"}), 404
    p = Path(__file__).parent / "data" / f"weekly-stats-{account_id}.json"
    if not p.exists() and account_id == "healthcare":
        p = Path(__file__).parent / "data" / "weekly-stats.json"
    if not p.exists():
        return jsonify({"error": "Not found"}), 503
    return jsonify(json.loads(p.read_text()))

# ── Account picker moved to templates/accounts.html ─────────────────────────────
_ACCOUNTS_HTML_UNUSED = """
<html lang="en">
<head>
  <meta charset="UTF-8"/>
  <meta name="viewport" content="width=device-width,initial-scale=1.0"/>
  <title>Company Signal Tracker</title>
  <link rel="icon" type="image/svg+xml" href="data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 32 32'%3E%3Crect width='32' height='32' rx='8' fill='%236366f1'/%3E%3Ccircle cx='16' cy='21' r='2.5' fill='%23fff'/%3E%3Cpath d='M10 15 Q16 9 22 15' stroke='%23fff' stroke-width='2' fill='none' stroke-linecap='round'/%3E%3Cpath d='M6 11 Q16 2 26 11' stroke='%23fff' stroke-width='2' fill='none' stroke-linecap='round' opacity='.55'/%3E%3C/svg%3E">
  <link rel="preconnect" href="https://fonts.googleapis.com"/>
  <link href="https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;500;600;700&display=swap" rel="stylesheet"/>
  <style>
    *,*::before,*::after{box-sizing:border-box;margin:0;padding:0}
    body{font-family:'Space Grotesk',sans-serif;color:#e2e8f0;
      min-height:100vh;display:flex;flex-direction:column;overflow-x:hidden;
      background:radial-gradient(ellipse 80% 50% at 20% 0%,rgba(99,102,241,.13) 0%,transparent 60%),
        radial-gradient(ellipse 60% 40% at 85% 60%,rgba(139,92,246,.10) 0%,transparent 55%),
        radial-gradient(ellipse 50% 60% at 50% 100%,rgba(30,27,75,.45) 0%,transparent 70%),
        linear-gradient(160deg,#080b18 0%,#0a0d1a 40%,#070912 100%)}
    .bg-grid{position:fixed;inset:0;z-index:0;pointer-events:none;
      background-image:radial-gradient(circle,rgba(99,102,241,.12) 1px,transparent 1px);
      background-size:36px 36px;
      mask-image:radial-gradient(ellipse 85% 85% at 50% 40%,black 30%,transparent 100%)}
    .bg-glow{position:fixed;border-radius:50%;filter:blur(130px);pointer-events:none;z-index:0;
      width:700px;height:700px;top:-200px;left:-150px;background:rgba(99,102,241,.08)}
    .topbar{position:relative;z-index:10;height:62px;padding:0 32px;
      display:flex;align-items:center;justify-content:space-between;
      background:rgba(7,9,16,.8);backdrop-filter:blur(16px);
      border-bottom:1px solid rgba(255,255,255,.05)}
    .tl{display:flex;align-items:center}
    .brand{display:flex;align-items:center;gap:10px;text-decoration:none}
    .brand-icon{width:34px;height:34px;border-radius:9px;
      background:linear-gradient(135deg,#6366f1,#8b5cf6);
      display:flex;align-items:center;justify-content:center;font-size:16px;
      box-shadow:0 0 14px rgba(99,102,241,.3)}
    .brand-name{font-size:15px;font-weight:700;color:#f1f5f9}
    .bc{display:flex;align-items:center;gap:8px;margin-left:18px;padding-left:18px;
      border-left:1px solid rgba(255,255,255,.07)}
    .bc a{font-size:13px;color:#2d3450;text-decoration:none;transition:color .15s}
    .bc a:hover{color:#64748b}
    .bc-sep{font-size:13px;color:#1a1d27}
    .bc-cur{font-size:13px;font-weight:600;color:#818cf8}
    .sign-out{font-size:12px;color:#3d4460;text-decoration:none;
      padding:6px 14px;border:1px solid rgba(255,255,255,.07);border-radius:8px;
      transition:all .15s}
    .sign-out:hover{color:#ef4444;border-color:rgba(239,68,68,.4)}
    .main{flex:1;position:relative;z-index:1;
      display:flex;flex-direction:column;align-items:center;padding:72px 24px 48px}
    .label{font-size:11px;font-weight:700;letter-spacing:.1em;text-transform:uppercase;
      color:#6366f1;margin-bottom:10px;display:flex;align-items:center;gap:8px}
    .label::before,.label::after{content:'';display:block;width:20px;height:1px;background:rgba(99,102,241,.4)}
    .heading{font-size:32px;font-weight:700;color:#f1f5f9;letter-spacing:-.02em;
      margin-bottom:6px;text-align:center}
    .sub{font-size:14px;color:#64748b;margin-bottom:52px}
    .grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(300px,360px));
      gap:20px;justify-content:center;width:100%;max-width:780px}
    .card{background:rgba(13,15,23,.9);border:1px solid rgba(255,255,255,.07);
      border-radius:22px;overflow:hidden;text-decoration:none;color:inherit;
      display:flex;flex-direction:column;
      transition:transform .22s cubic-bezier(.34,1.56,.64,1),box-shadow .22s,border-color .2s}
    .card:hover{transform:translateY(-5px);
      box-shadow:0 24px 64px rgba(0,0,0,.55),0 0 0 1px var(--glow)}
    .card-band{height:3px;background:var(--accent)}
    .card-thumb{height:110px;background:var(--thumb);position:relative;
      display:flex;align-items:center;justify-content:center;overflow:hidden}
    .card-thumb-icon{font-size:44px;opacity:.45;filter:drop-shadow(0 0 20px rgba(255,255,255,.15))}
    .card-thumb::after{content:'';position:absolute;inset:0;
      background:linear-gradient(to bottom,transparent 30%,rgba(13,15,23,.95) 100%)}
    .card-badge{position:absolute;top:10px;right:10px;z-index:1;
      font-size:9px;font-weight:700;letter-spacing:.07em;text-transform:uppercase;
      padding:3px 9px;border-radius:999px;display:flex;align-items:center;gap:4px;
      background:rgba(16,185,129,.15);border:1px solid rgba(16,185,129,.3);color:#34d399}
    .badge-dot{width:5px;height:5px;border-radius:50%;background:currentColor;
      animation:bpulse 2s infinite}
    @keyframes bpulse{0%,100%{box-shadow:0 0 0 0 rgba(52,211,153,.5)}
      50%{box-shadow:0 0 0 3px rgba(52,211,153,0)}}
    .card-body{padding:20px 24px 22px;flex:1;display:flex;flex-direction:column}
    .card-name{font-size:20px;font-weight:700;color:#f1f5f9;letter-spacing:-.01em;margin-bottom:8px}
    .card-desc{font-size:13px;color:#94a3b8;line-height:1.65;flex:1;margin-bottom:20px}
    .card-footer{display:flex;align-items:center;justify-content:space-between;
      border-top:1px solid rgba(255,255,255,.07);padding-top:16px}
    .stat{font-size:12px;color:#64748b}
    .stat span{color:var(--accent-text);font-weight:600}
    .arrow{font-size:16px;color:var(--accent-text);opacity:0;transition:opacity .15s,transform .15s}
    .card:hover .arrow{opacity:1;transform:translateX(3px)}
    .foot{margin-top:48px;font-size:12px;color:#13151f}
  </style>
</head>
<body>
  <div class="bg-grid"></div>
  <div class="bg-glow"></div>
  <div class="topbar">
    <div class="tl">
      <a href="/hub" class="brand">
        <div class="brand-icon">📡</div>
        <span class="brand-name">Platform</span>
      </a>
      <div class="bc">
        <a href="/hub">Hub</a><span class="bc-sep">›</span>
        <a href="/ppc">PPC</a><span class="bc-sep">›</span>
        <span class="bc-cur">Signal Tracker</span>
      </div>
    </div>
    <a href="/logout" class="sign-out">Sign out</a>
  </div>
  <div class="main">
    <div class="label">Company Intelligence</div>
    <h1 class="heading">Company Signal Tracker</h1>
    <p class="sub">Choose a company list to open the dashboard</p>
    <div class="grid">{account_cards}</div>
    <p class="foot">Position2 · Internal use only</p>
  </div>
</body>
</html>"""


def _build_account_card(account_id, cfg):
    path = cfg["dashboard"]
    accent = cfg["accent"]
    # derive thumb gradient from accent colour
    thumb_map = {"#3b82f6": "linear-gradient(135deg,#172554,#1e3a8a)",
                 "#8b5cf6": "linear-gradient(135deg,#2e1065,#1e1b4b)"}
    thumb = thumb_map.get(accent, f"linear-gradient(135deg,#0d0f17,#1a1d27)")
    if path.exists():
        count = _read_company_count(path)
        refreshed = _read_last_refreshed(path)
        return (
            f'<a class="card" href="/signal-tracker/{account_id}" '
            f'style="--accent:{accent};--glow:rgba(99,102,241,.25);'
            f'--thumb:{thumb};--accent-text:{accent}">'
            f'<div class="card-band"></div>'
            f'<div class="card-thumb"><div class="card-thumb-icon">{cfg["icon"]}</div>'
            f'<div class="card-badge"><span class="badge-dot"></span>Live</div></div>'
            f'<div class="card-body">'
            f'<div class="card-name">{cfg["name"]}</div>'
            f'<div class="card-desc">{cfg["description"]}</div>'
            f'<div class="card-footer">'
            f'<div class="stat"><span>{count}</span> companies</div>'
            f'<span class="arrow">→</span></div></div></a>'
        )
    return (
        f'<div class="card" style="--accent:{accent};--glow:rgba(99,102,241,.15);'
        f'--thumb:{thumb};--accent-text:{accent};opacity:.5;cursor:default">'
        f'<div class="card-band"></div>'
        f'<div class="card-thumb"><div class="card-thumb-icon">{cfg["icon"]}</div></div>'
        f'<div class="card-body">'
        f'<div class="card-name">{cfg["name"]}</div>'
        f'<div class="card-desc">{cfg["description"]}</div>'
        f'<div class="card-footer">'
        f'<div class="stat" style="color:#f59e0b">Not generated yet</div>'
        f'</div></div></div>'
    )


def _read_last_refreshed(path: Path) -> str:
    try:
        mtime = path.stat().st_mtime
        dt = datetime.fromtimestamp(mtime, tz=IST)
        d = dt.strftime("%d %b %Y").lstrip("0")
        t = dt.strftime("%I:%M %p").lstrip("0")
        return f"{d}, {t} IST"
    except Exception:
        return "unknown"


def _read_company_count(path: Path) -> str:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
        idx = text.find('"total_companies":')
        if idx == -1:
            return "—"
        snippet = text[idx + 18:idx + 28].strip().split(",")[0].strip()
        return snippet if snippet.isdigit() else "—"
    except Exception:
        return "—"


# ── Shared Sheets helper ──────────────────────────────────────────────────────

def _sheets_service():
    """Return an authenticated Google Sheets service, or raise on failure."""
    import json as _j
    from google.oauth2 import service_account
    from googleapiclient.discovery import build

    sa_str = os.environ.get("GOOGLE_SA_JSON", "")
    if not sa_str:
        raise RuntimeError("GOOGLE_SA_JSON env var not set")
    sa_info = _j.loads(sa_str)
    creds = service_account.Credentials.from_service_account_info(
        sa_info, scopes=["https://www.googleapis.com/auth/spreadsheets.readonly"])
    return build("sheets", "v4", credentials=creds, cache_discovery=False)


# ── Chatbot data functions ────────────────────────────────────────────────────

def _chatbot_get_anonymous_visitors(date_from=None, date_to=None, company=None,
                                     seniority=None, industry=None, limit=20):
    """Fetch targeted anonymous visitor rows for the chatbot."""
    try:
        svc = _sheets_service()
        r = svc.spreadsheets().values().get(
            spreadsheetId=ANON_VISITORS_SHEET_ID,
            range="People Enriched!A:L"
        ).execute()
        rows = r.get("values", [])
        if not rows:
            return {"status": "empty", "message": "Sheet returned no data", "people": []}

        # Header-based mapping — robust to column reordering
        headers = [h.strip().lower() for h in rows[0]]

        def _h(row, *names):
            """Get the first matching column value by header name."""
            for name in names:
                for i, h in enumerate(headers):
                    if name in h and i < len(row):
                        return row[i]
            return ""

        people = []
        for row in rows[1:]:
            name = _h(row, "name", "full name")
            if not name or name.strip().lower() in ("", "unavailable", "n/a"):
                continue
            time_str = _h(row, "time", "date", "timestamp", "visited", "last seen")
            industry_raw = _h(row, "industry", "sector")
            people.append({
                "name":     name,
                "title":    _h(row, "title", "job title", "position", "role"),
                "email":    _h(row, "email"),
                "company":  _h(row, "company", "organization", "employer", "account"),
                "location": _h(row, "location", "city", "region", "country"),
                "pages":    _h(row, "pages", "page", "url", "viewed"),
                "date":     time_str[:10] if time_str else "",
                "industry": _clean_industry(industry_raw),
                "website":  _h(row, "website", "domain", "web", "url"),
                "time_raw": time_str,
            })

        # Newest first
        people.sort(key=lambda x: x.get("time_raw", ""), reverse=True)

        total_before_filter = len(people)

        if date_from:
            people = [p for p in people if p["date"] >= date_from]
        if date_to:
            people = [p for p in people if p["date"] <= date_to]
        if company:
            c = company.lower()
            people = [p for p in people
                      if c in p.get("company", "").lower()
                      or c in p.get("website", "").lower()]
        if industry:
            people = [p for p in people
                      if industry.lower() in p.get("industry", "").lower()]
        if seniority:
            s = seniority.lower()
            _seniority_map = {
                "c-suite":   ["ceo", "cmo", "coo", "cto", "cfo", "cro", "cpo", "ciso", "chief"],
                "vp":        ["vp", "vice president"],
                "director":  ["director"],
                "manager":   ["manager"],
                "president": ["president"],
            }
            keywords = _seniority_map.get(s, [s])
            people = [p for p in people
                      if any(kw in p.get("title", "").lower() for kw in keywords)]

        result = people[:limit]
        # Industry breakdown
        industry_counts = dict(Counter(p["industry"] for p in people if p["industry"]).most_common(5))
        return {
            "status": "ok",
            "total_in_sheet": total_before_filter,
            "total_matching_filters": len(people),
            "returned": len(result),
            "top_industries": industry_counts,
            "people": [
                {
                    "name":         p["name"],
                    "title":        p["title"],
                    "company":      p["company"],
                    "industry":     p["industry"],
                    "location":     p["location"],
                    "date_visited": p["date"],
                    "pages_viewed": p["pages"],
                }
                for p in result
            ],
        }
    except Exception as e:
        return {"status": "error", "error": str(e),
                "hint": "Check that GOOGLE_SA_JSON is set and the sheet is accessible."}


def _chatbot_get_signal_tracker(account="healthcare", signal_type=None,
                                 company=None, severity=None, limit=20):
    """Query Signal Tracker SQLite for buying signals."""
    db_map = {
        "healthcare": Path(__file__).parent / "data" / "tracker.db",
        "csg":        Path(__file__).parent / "data" / "tracker_csg_v2.db",
    }
    db_path = db_map.get(account, db_map["healthcare"])
    if not db_path.exists():
        return {"error": f"Database not found for account '{account}'"}

    try:
        import sqlite3 as _sql
        conn = _sql.connect(str(db_path))
        conn.row_factory = _sql.Row

        conditions = ["a.dry_run = 0"]
        params: list = []
        if signal_type:
            conditions.append("a.signal_type = ?")
            params.append(signal_type)
        if severity:
            conditions.append("a.severity = ?")
            params.append(severity.upper())
        if company:
            conditions.append("c.name LIKE ?")
            params.append(f"%{company}%")

        where = " AND ".join(conditions)
        query = f"""
            SELECT c.name, c.domain, c.industry, c.city, c.state,
                   a.signal_type, a.signal_detail, a.severity, a.signal_date
            FROM alerts_sent a
            JOIN companies c ON a.apollo_id = c.apollo_id
            WHERE {where}
            ORDER BY a.signal_date DESC
            LIMIT ?
        """
        params.append(limit)
        rows = conn.execute(query, params).fetchall()
        conn.close()

        signals = [
            {
                "company":       r["name"],
                "domain":        r["domain"],
                "industry":      r["industry"],
                "location":      f"{r['city']}, {r['state']}".strip(", "),
                "signal_type":   r["signal_type"],
                "signal_detail": r["signal_detail"],
                "severity":      r["severity"],
                "signal_date":   r["signal_date"],
            }
            for r in rows
        ]
        return {"account": account, "total_returned": len(signals), "signals": signals}

    except Exception as e:
        return {"error": str(e)}


# ── OpenAI chatbot definitions ────────────────────────────────────────────────

CHATBOT_FUNCTIONS = [
    {
        "name": "get_anonymous_visitors",
        "description": (
            "Get people who visited position2.com — identified and enriched via Apollo. "
            "Use for: who visited, how many, which companies, seniority levels, industry breakdown, "
            "recent visitors, visitors in a date range."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "date_from": {"type": "string", "description": "Start date YYYY-MM-DD"},
                "date_to":   {"type": "string", "description": "End date YYYY-MM-DD"},
                "company":   {"type": "string", "description": "Filter by company name or website domain"},
                "seniority": {
                    "type": "string",
                    "description": "Filter by seniority: 'c-suite', 'vp', 'director', 'manager', or any title keyword",
                },
                "industry":  {"type": "string", "description": "Filter by industry keyword"},
                "limit":     {"type": "integer", "description": "Max results (default 20)"},
            },
        },
    },
    {
        "name": "get_signal_tracker",
        "description": (
            "Get buying signals from the Signal Tracker — companies showing funding rounds, "
            "C-suite changes, M&A, news mentions, or IPO signals. "
            "Use for: prospect intelligence, hot accounts, recent high signals, outbound prioritization."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "account": {
                    "type": "string",
                    "enum": ["healthcare", "csg"],
                    "description": "Which tracker — 'healthcare' (1,251 companies) or 'csg' (294 companies). Default: healthcare",
                },
                "signal_type": {
                    "type": "string",
                    "description": (
                        "Filter by signal type: 'Funding Round', 'C-Suite Join', 'C-Suite Exit', "
                        "'Acquisition / M&A', 'News Mention', 'IPO Signal'"
                    ),
                },
                "company":  {"type": "string", "description": "Filter by company name"},
                "severity": {"type": "string", "enum": ["HIGH", "LOW"], "description": "HIGH = funding/C-suite/M&A; LOW = news"},
                "limit":    {"type": "integer", "description": "Max results (default 20)"},
            },
        },
    },
    {
        "name": "get_ad_intelligence_data",
        "description": (
            "Fetch competitor ad data from Ad Intelligence. "
            "Competitors tracked: Inspire Aesthetics, Dr. Dana MD, Sono Bello. "
            "Use for: what ads competitors are running, CTAs, ad formats, keywords targeted, "
            "messaging angles, active vs inactive ads, when ads were last seen."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "competitor": {"type": "string", "description": "Competitor name or domain (e.g. 'Inspire Aesthetics', 'sonobello')"},
                "ad_format":  {"type": "string", "enum": ["image", "text", "video"], "description": "Ad format filter"},
                "status":     {"type": "string", "enum": ["active", "inactive"], "description": "Ad status filter"},
                "keyword":    {"type": "string", "description": "Search word in headline/description/keywords"},
                "limit":      {"type": "integer", "description": "Max ads to return (default 20)"},
            },
        },
    },
]

_PPC_CTX_CACHE: dict = {"data": None, "ts": 0.0}
_PPC_CTX_TTL = 60    # seconds — refresh every 60s; keeps data fresh without hammering APIs


def _build_ppc_context() -> str:
    """
    Fetch ALL data from every source — no row limits.
    Cached for _PPC_CTX_TTL seconds so repeated chat messages are instant.
    """
    import time as _time
    now = _time.time()
    if _PPC_CTX_CACHE["data"] and now - _PPC_CTX_CACHE["ts"] < _PPC_CTX_TTL:
        return _PPC_CTX_CACHE["data"]

    parts = []

    # ── 1. Anonymous Visitors — ALL rows from BOTH tabs ────────────────────
    try:
        svc = _sheets_service()

        # ---- Tab 1: People Enriched (individual visitors) ----------------
        r_people = svc.spreadsheets().values().get(
            spreadsheetId=ANON_VISITORS_SHEET_ID,
            range="People Enriched!A:L"
        ).execute()
        people_rows = r_people.get("values", [])

        # Read header row to map columns by name
        headers = [h.strip().lower() for h in (people_rows[0] if people_rows else [])]

        def _hv(row, *names):
            for name in names:
                for i, h in enumerate(headers):
                    if name in h and i < len(row):
                        v = row[i].strip()
                        if v and v.lower() not in ("", "unavailable", "n/a", "none"):
                            return v
            return ""

        people_out = []
        for row in people_rows[1:]:
            name = _hv(row, "name", "full name")
            if not name:
                continue
            time_str = _hv(row, "time", "date", "timestamp", "identified", "visited", "seen", "first")
            people_out.append({
                "name":     name,
                "title":    _hv(row, "title", "job title", "position", "role"),
                "email":    _hv(row, "email"),
                "location": _hv(row, "location", "city", "country", "region"),
                "pages":    _hv(row, "pages", "page", "url"),
                "date":     time_str[:10] if time_str else "",
                "industry": _clean_industry(_hv(row, "industry", "sector", "vertical")),
                "website":  _hv(row, "website", "domain", "company website", "web"),
                "time_raw": time_str,
            })
        people_out.sort(key=lambda x: x.get("time_raw", ""), reverse=True)

        # ---- Tab 2: Visitors By Company (company-level data) --------------
        r_comp = svc.spreadsheets().values().get(
            spreadsheetId=ANON_VISITORS_SHEET_ID,
            range="Visitors By Company!A:J"
        ).execute()
        comp_rows = r_comp.get("values", [])
        comp_hdrs = [h.strip().lower() for h in (comp_rows[0] if comp_rows else [])]

        def _cv(row, *names):
            for name in names:
                for i, h in enumerate(comp_hdrs):
                    if name in h and i < len(row):
                        v = row[i].strip()
                        if v and v.lower() not in ("", "unavailable", "n/a", "none"):
                            return v
            return ""

        companies_out = []
        for row in comp_rows[1:]:
            co = _cv(row, "company", "name", "organization") or (row[0].strip() if row else "")
            if not co:
                continue
            companies_out.append({
                "company":   co,
                "website":   _cv(row, "website", "domain", "url") or (row[2].strip() if len(row) > 2 else ""),
                "location":  " ".join(filter(None, [_cv(row, "city"), _cv(row, "state"), _cv(row, "country")])),
                "industry":  _clean_industry(_cv(row, "industry", "sector")),
                "employees": _cv(row, "employee", "size", "headcount"),
                "revenue":   _cv(row, "revenue", "arr", "mrr"),
            })

        # Company lines — fields explicitly labelled so GPT never guesses column order
        c_lines = []
        for i, c in enumerate(companies_out, 1):
            c_lines.append(
                f"{i}. Company={c['company']} | Website={c['website']} | "
                f"Industry={c['industry']} | Location={c['location']} | "
                f"Employees={c['employees']} | Revenue={c['revenue']}"
            )

        # People lines — completely separate block with different field set
        p_lines = []
        for i, p in enumerate(people_out, 1):
            p_lines.append(
                f"{i}. Name={p['name']} | Title={p['title']} | "
                f"CompanyWebsite={p['website']} | Industry={p['industry']} | "
                f"Location={p['location']} | DateVisited={p['date']}"
            )

        industry_counts = dict(Counter(p["industry"] for p in people_out if p["industry"]).most_common(8))

        # COMPANIES block comes FIRST so GPT reads it first for "company" queries
        parts.append(
            f"=== VISITOR DATA ===\n"
            f"Summary: {len(people_out)} individual visitors from {len(companies_out)} unique companies\n"
            f"Top industries: {industry_counts}\n\n"
            f"--- SECTION A: COMPANIES THAT VISITED ({len(companies_out)} unique companies) ---\n"
            f"USE THIS SECTION when asked about COMPANIES. Columns: Company, Website, Industry, Location, Employees, Revenue\n"
            + "\n".join(c_lines)
            + f"\n\n--- SECTION B: INDIVIDUAL VISITORS ({len(people_out)} people, newest first) ---\n"
            f"USE THIS SECTION when asked about VISITORS or PEOPLE. Columns: Name, Title, CompanyWebsite, Industry, Location, DateVisited\n"
            + "\n".join(p_lines)
        )

    except Exception as e:
        parts.append(f"=== ANONYMOUS VISITORS ===\n⚠ Could not fetch: {e}")

    # ── 2. Signal Tracker — ALL signals, no limit ─────────────────────────
    try:
        import sqlite3 as _sql

        db_path = Path(__file__).parent / "data" / "tracker.db"
        if not db_path.exists():
            parts.append("=== SIGNAL TRACKER ===\n⚠ Database not on Railway — commit data/tracker.db to git")
        else:
            conn = _sql.connect(str(db_path))
            conn.row_factory = _sql.Row
            try:
                all_sigs = conn.execute("""
                    SELECT c.name, c.domain, c.industry, c.city, c.state,
                           a.signal_type, a.signal_detail, a.severity, a.signal_date
                    FROM alerts_sent a
                    JOIN companies c ON a.apollo_id = c.apollo_id
                    WHERE a.dry_run = 0
                    ORDER BY a.signal_date DESC
                """).fetchall()
            finally:
                conn.close()

            sig_counts = dict(Counter(r["signal_type"] for r in all_sigs).most_common())
            comp_count = len({r["name"] for r in all_sigs})

            sig_lines = []
            for r in all_sigs:
                date   = (r["signal_date"] or "")[:16]
                detail = (r["signal_detail"] or "")[:120]
                city   = r["city"] or ""
                state  = r["state"] or ""
                loc    = ", ".join(filter(None, [city, state]))
                sig_lines.append(
                    f"• {r['name']} | {r['industry']} | {loc} | "
                    f"{r['signal_type']} [{r['severity']}] | {date} | {detail}"
                )

            parts.append(
                f"=== SIGNAL TRACKER (Healthcare — 1,251 companies monitored) ===\n"
                f"Total signals: {len(all_sigs)} across {comp_count} companies\n"
                f"By type: {sig_counts}\n\n"
                f"--- ALL SIGNALS (newest first) ---\n"
                + "\n".join(sig_lines)
            )

    except Exception as e:
        parts.append(f"=== SIGNAL TRACKER ===\n⚠ Could not fetch: {e}")

    # ── 3. Ad Intelligence — ALL ads ─────────────────────────────────────
    try:
        a = get_ad_intelligence_data(limit=5000)   # effectively unlimited
        if a.get("status") == "ok":
            ad_lines = []
            for ad in a.get("ads", []):
                ad_lines.append(
                    f"• {ad['competitor']} | {ad['format']} | {ad['status']} | "
                    f"headline: '{ad['headline']}' | CTA: {ad['cta']} | "
                    f"keywords: {ad['keywords'][:80]} | "
                    f"angle: {ad['messaging_angle'][:60]} | "
                    f"first: {ad['first_shown']} | last: {ad['last_shown']}"
                )
            parts.append(
                f"=== AD INTELLIGENCE ===\n"
                f"Total ads tracked: {a['total_in_sheet']}\n"
                f"By competitor: {a['by_competitor']}\n"
                f"By format: {a['by_format']}\n"
                f"By status: {a['by_status']}\n"
                f"Top CTAs: {a['top_ctas']}\n"
                f"Top keywords: {a['top_keywords']}\n\n"
                f"--- ALL ADS ---\n"
                + "\n".join(ad_lines)
            )
        else:
            fix = a.get("fix", "")
            parts.append(
                f"=== AD INTELLIGENCE ===\n⚠ {a.get('error','Error')}\n"
                + (f"ACTION: {fix}" if fix else "")
            )
    except Exception as e:
        parts.append(f"=== AD INTELLIGENCE ===\n⚠ Could not fetch: {e}")

    ctx = "\n\n" + "\n\n".join(parts)
    _PPC_CTX_CACHE["data"] = ctx
    _PPC_CTX_CACHE["ts"] = now
    return ctx


@app.route("/api/ppc-chat", methods=["POST"])
@login_required
def ppc_chat():
    """PPC AI assistant — context injection."""
    from openai import OpenAI
    api_key = os.environ.get("OPENAI_API_KEY", "")
    if not api_key:
        return jsonify({"answer": "⚠️ Add `OPENAI_API_KEY` to Railway Variables."}), 200

    oai = OpenAI(api_key=api_key)

    body         = request.json or {}
    user_message = body.get("message", "").strip()
    history      = body.get("history", [])[-12:]
    memories     = body.get("memories", [])
    source_text  = body.get("source_text", "")   # previous AI response to reformat
    export_fmt   = body.get("export_format", "")  # "csv", "excel", "table", "json", etc.

    # ── Attached file (optional) ────────────────────────────────────────────
    file_name     = body.get("file_name", "")
    file_content  = body.get("file_content", "")
    file_is_image = body.get("file_is_image", False)
    file_mime     = body.get("file_mime", "image/png")
    file_base64   = body.get("file_base64", "")
    file_truncated= body.get("file_truncated", False)
    has_file      = bool(file_name)

    if not user_message and not has_file:
        return jsonify({"answer": "Please ask a question."}), 200

    if not user_message and has_file:
        user_message = f"Please analyse the attached file '{file_name}' and summarise the key information."

    # ── Detect format keyword in the user's message itself ───────────────────
    _fmt_map = {
        r'\bexcel\b|\bxlsx\b':                              'excel',
        r'\bcsv\b':                                         'csv',
        r'\bjson\b':                                        'json',
        r'\btable format\b|\bin a table\b|\bspreadsheet\b': 'table',
        r'\bbullet\b|\blist format\b':                      'bullet',
    }
    if not export_fmt:
        for pattern, fmt in _fmt_map.items():
            if re.search(pattern, user_message, re.I):
                export_fmt = fmt
                break

    # ── FORMAT / EXPORT REQUEST — handle separately, no PPC context needed ──
    # Triggered when user asks to reformat a *previous* response
    if export_fmt and source_text:
        fmt_instructions = {
            "csv":   "Convert the data to clean CSV with a header row. Use comma separators. Output ONLY the CSV — no explanation, no markdown, no code block.",
            "excel": "Convert the data to clean CSV with a header row (Excel-compatible). Use comma separators. Output ONLY the CSV data — no explanation, no markdown, no code block.",
            "table": "Format the data as a clean markdown table with aligned columns and a header row. Output ONLY the table.",
            "json":  "Convert the data to valid JSON array of objects. Output ONLY the JSON — no explanation, no markdown.",
            "bullet":"Reformat the data as a clean bulleted list. Output ONLY the list.",
        }
        instruction = fmt_instructions.get(export_fmt, f"Reformat the data as {export_fmt}. Output ONLY the reformatted data.")
        reformat_messages = [
            {"role": "system", "content":
             "You are a data formatter. Your only job is to reformat data exactly as instructed. "
             "Never add explanations, apologies, or commentary. Output ONLY the requested format."},
            {"role": "user", "content": f"{instruction}\n\nDATA TO REFORMAT:\n{source_text}"},
        ]
        try:
            formatted, _m = _kairo_completion(oai, reformat_messages, 2000, temperature=0)
            is_csv = export_fmt in ("csv", "excel")
            return jsonify({"answer": formatted, "is_export": True,
                            "export_format": export_fmt, "is_csv": is_csv})
        except Exception as e:
            return jsonify({"answer": f"Export failed: {str(e)}"}), 200

    # ── Pre-fetch all live data (cached 4 min) ─────────────────────────────
    ppc_context = _build_ppc_context()

    now_ist    = datetime.now(IST)
    today      = now_ist.strftime("%Y-%m-%d")
    week_start = (now_ist - timedelta(days=now_ist.weekday())).strftime("%Y-%m-%d")

    # Build format instruction if a format was requested in this message
    fmt_instruction = ""
    if export_fmt:
        fmt_map = {
            "excel": "Return the data as clean CSV (Excel-compatible) with a header row. Only output the CSV rows — no prose, no markdown fences.",
            "csv":   "Return the data as clean CSV with a header row. Only output the CSV rows — no prose, no markdown fences.",
            "json":  "Return the data as a valid JSON array of objects. Only output the JSON.",
            "table": "Return the data formatted as a clean markdown table with headers.",
        }
        fmt_instruction = f"\n\nOUTPUT FORMAT REQUIRED: {fmt_map.get(export_fmt, f'Format the output as {export_fmt}.')}\nDo NOT include any explanation before or after the data."

    system_prompt = f"""You are the PPC Intelligence Assistant for Position2, a B2B marketing agency.
You are highly intelligent, direct, and always give complete answers in one response — no follow-up questions.

TODAY: {today} | THIS WEEK: {week_start} to {today} | YESTERDAY: {(now_ist - timedelta(days=1)).strftime('%Y-%m-%d')}
"Last N" = first N rows of the relevant list (data is newest-first).

INSTRUCTIONS:
- Answer every question fully using the live data below. Never say "I can't access" when data is provided.
- Never ask for clarification when the request is clear. Deliver the answer immediately.
- "Excel format", "CSV", "table", "JSON" = format the output that way. Nothing to do with any Google Sheet.
- If a data section shows ⚠ Error, say that source is unavailable but answer from what's available.
- Be analytical: bold **key numbers**, use bullets for lists, lead with the most useful insight.
- For general PPC/marketing questions, answer from knowledge directly.

DATA SECTION RULES — NEVER MIX THESE:
- Asked about COMPANIES → use SECTION A only. Columns: Company Name, Website, Industry, Location, Employees, Revenue. Never include individual people names.
- Asked about VISITORS/PEOPLE → use SECTION B only. Columns: Name, Title, Company Website, Industry, Location, Date Visited.
- "last 10 companies" = first 10 rows of SECTION A. "last 10 visitors" = first 10 rows of SECTION B.

CSV/EXCEL EXPORT RULES:
- Output ONLY the CSV rows. No intro text, no explanation, no markdown fences, no code blocks.
- Use meaningful headers: "Company Name", "Website", "Industry", "Location", "Employees", "Revenue" — never "field1", "field2".
- Include ONLY the columns that make sense for the query (e.g. company query = 6 columns, no extra).
- Replace em-dashes (—) with a hyphen or leave blank. Quote values that contain commas.{fmt_instruction}

══════════════════════════ LIVE DATA ══════════════════════════
{ppc_context}
═══════════════════════════════════════════════════════════════
"""
    if memories:
        system_prompt += "\n\nSAVED MEMORIES (always apply these):\n" + \
                         "\n".join(f"• {m}" for m in memories[:30])

    messages = [{"role": "system", "content": system_prompt}]
    messages += history

    # ── Build user turn — plain text or multimodal (image) ─────────────────
    if file_is_image and file_base64:
        # Vision: send image alongside the question
        trunc_note = " (image sent in full)"
        messages.append({
            "role": "user",
            "content": [
                {"type": "text",
                 "text": f"I've attached an image: '{file_name}'\n{user_message}"},
                {"type": "image_url",
                 "image_url": {"url": f"data:{file_mime};base64,{file_base64}", "detail": "high"}},
            ],
        })
    elif file_content:
        # Text file: prepend content clearly labelled
        trunc_note = f"\n\n⚠️ File was large — showing first {_MAX_TEXT_CHARS:,} characters." if file_truncated else ""
        file_block = (
            f"📎 ATTACHED FILE: {file_name}{trunc_note}\n"
            f"{'─'*60}\n"
            f"{file_content}\n"
            f"{'─'*60}\n\n"
            f"User question about this file: {user_message}"
        )
        messages.append({"role": "user", "content": file_block})
    else:
        messages.append({"role": "user", "content": user_message})

    try:
        answer, _m = _kairo_completion(oai, messages, 2000, temperature=0.1)
        return jsonify({
            "answer":          answer,
            "detected_format": export_fmt or "",
            "is_export":       bool(export_fmt and not source_text),
            "is_csv":          export_fmt in ("csv", "excel"),
        })
    except Exception as e:
        log.warning("PPC chat error: %s", e)
        return jsonify({"answer": f"Something went wrong: {str(e)}"}), 200


# ── File extraction helpers ───────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> str:
    import pdfplumber, io
    parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                parts.append(f"[Page {i}]\n{text.strip()}")
            # Extract tables too
            for table in page.extract_tables():
                rows = [" | ".join(str(c) if c else "" for c in row) for row in table if any(c for c in row)]
                if rows:
                    parts.append("\n".join(rows))
    return "\n\n".join(parts) or "(No text could be extracted from this PDF)"


def _extract_docx(data: bytes) -> str:
    import docx, io
    doc = docx.Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            style = para.style.name if para.style else ""
            prefix = "# " if "Heading 1" in style else "## " if "Heading" in style else ""
            parts.append(prefix + para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = " | ".join(c.text.strip() for c in row.cells)
            if cells.strip():
                parts.append(cells)
    return "\n".join(parts) or "(No text found in document)"


def _extract_xlsx(data: bytes) -> str:
    import openpyxl, io
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            vals = [str(v).strip() if v is not None else "" for v in row]
            if any(v for v in vals):
                parts.append(" | ".join(vals))
    return "\n".join(parts) or "(No data found in spreadsheet)"


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    import io
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_parts))
    return "\n\n".join(parts) or "(No text found in presentation)"


_IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "bmp"}
_IMAGE_MIME  = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                "gif": "image/gif",  "webp": "image/webp", "bmp": "image/bmp"}
_MAX_FILE_BYTES  = 20 * 1024 * 1024   # 20 MB
_MAX_TEXT_CHARS  = 40_000             # truncate extracted text at 40k chars


@app.route("/api/ppc-upload", methods=["POST"])
@login_required
def ppc_upload():
    """Parse an uploaded file and return its extracted content for the chatbot."""
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    f        = request.files["file"]
    filename = f.filename or "upload"
    ext      = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    data     = f.read()

    if len(data) > _MAX_FILE_BYTES:
        return jsonify({"error": f"File too large — max {_MAX_FILE_BYTES // 1024 // 1024} MB"}), 400
    if not ext:
        return jsonify({"error": "File has no extension — cannot detect type"}), 400

    try:
        # ── Images → return base64 for vision API ──────────────────────────
        if ext in _IMAGE_EXTS:
            import base64 as _b64
            b64  = _b64.b64encode(data).decode()
            mime = _IMAGE_MIME.get(ext, "image/png")
            return jsonify({"is_image": True, "name": filename, "mime": mime, "base64": b64})

        # ── Text-based files → extract and return text ─────────────────────
        if ext == "pdf":
            content = _extract_pdf(data)
        elif ext in ("docx", "doc"):
            content = _extract_docx(data)
        elif ext in ("xlsx", "xls"):
            content = _extract_xlsx(data)
        elif ext in ("pptx", "ppt"):
            content = _extract_pptx(data)
        elif ext in ("csv", "txt", "md", "json", "xml", "html", "htm"):
            content = data.decode("utf-8", errors="replace")
        else:
            return jsonify({"error": f"Unsupported file type .{ext}. Supported: PDF, DOCX, XLSX, PPTX, CSV, TXT, PNG, JPG, and more."}), 400

        # Truncate if huge
        truncated = False
        if len(content) > _MAX_TEXT_CHARS:
            content   = content[:_MAX_TEXT_CHARS]
            truncated = True

        return jsonify({
            "is_image": False,
            "name":      filename,
            "content":   content,
            "chars":     len(content),
            "truncated": truncated,
        })

    except Exception as e:
        log.warning("ppc_upload error for %s: %s", filename, e)
        return jsonify({"error": f"Could not parse '{filename}': {str(e)}"}), 500



# ── Signal Tracker Insights API ──────────────────────────────────────────────

_REVENUE_KEYS = {"est_value", "opportunity", "revenue_impact", "pipeline_estimate",
                 "estimated_value", "pipeline_value", "deal_size", "contract_value"}

def _responses_web_search(oai, model, input_msgs, max_tokens):
    """Call the Responses API with web search, trying both known tool-type names
    ('web_search' and the older 'web_search_preview'). Returns (text, True) on
    success or (None, False) if web search is unavailable on this SDK/model."""
    for _tt in ("web_search", "web_search_preview"):
        try:
            resp = oai.responses.create(
                model=model, tools=[{"type": _tt}], input=input_msgs, max_output_tokens=max_tokens)
            txt = (getattr(resp, "output_text", "") or "").strip()
            if txt:
                return txt, True
        except Exception as we:
            log.warning("web search via '%s' unavailable: %s", _tt, we)
    return None, False


def _kairo_model_chain():
    """Strongest-first model chain: OPENAI_INSIGHTS_MODEL > gpt-5.4 (ChatGPT 5.4) > OPENAI_MODEL/gpt-4o-mini."""
    chain = []
    for m in (os.environ.get("OPENAI_INSIGHTS_MODEL"), "gpt-5.4",
              os.environ.get("OPENAI_MODEL", "gpt-4o-mini")):
        if m and m not in chain:
            chain.append(m)
    return chain


def _kairo_completion(oai, messages, max_tokens, temperature=None):
    """Plain-text chat completion on the primary Kairo model (GPT-5.4) with graceful
    fallback down the model chain; retries without temperature if a model rejects it.
    Returns (text, model_used)."""
    last_err = None
    for model in _kairo_model_chain():
        attempts = [{"temperature": temperature}] if temperature is not None else []
        attempts.append({})
        for kw in attempts:
            try:
                resp = oai.chat.completions.create(
                    model=model, messages=messages,
                    max_completion_tokens=max_tokens, **kw)
                txt = (resp.choices[0].message.content or "").strip()
                if txt:
                    return txt, model
            except Exception as e:
                last_err = e
                log.warning("kairo: completion on '%s' (%s) failed: %s", model, kw, e)
    raise last_err if last_err else RuntimeError("no usable OpenAI model")


def _kairo_chat_json(oai, messages, max_tokens):
    """Chat completion in strict JSON mode, trying the strongest model first.
    Returns (raw_text, model_used)."""
    last_err = None
    for model in _kairo_model_chain():
        try:
            resp = oai.chat.completions.create(
                model=model, messages=messages,
                max_completion_tokens=max_tokens,
                response_format={"type": "json_object"})
            txt = (resp.choices[0].message.content or "").strip()
            if txt:
                return txt, model
        except Exception as e:
            last_err = e
            log.warning("kairo: model '%s' failed, trying next: %s", model, e)
    raise last_err if last_err else RuntimeError("no usable OpenAI model")


def _strip_revenue_fields(obj):
    """Recursively remove all revenue / pipeline-value fields from GPT output."""
    if isinstance(obj, dict):
        return {k: _strip_revenue_fields(v) for k, v in obj.items() if k not in _REVENUE_KEYS}
    if isinstance(obj, list):
        return [_strip_revenue_fields(x) for x in obj]
    return obj

@app.route("/api/insights-meta/<account_id>")
@login_required
def insights_meta(account_id):
    import sqlite3
    from pathlib import Path
    db_map = {"healthcare": Path(__file__).parent/"data"/"tracker.db",
              "csg":        Path(__file__).parent/"data"/"tracker_csg_v2.db"}
    db_path = db_map.get(account_id)
    if not db_path or not db_path.exists():
        return jsonify({"error": "Unknown account"}), 404
    try:
        conn = sqlite3.connect(str(db_path)); conn.row_factory = sqlite3.Row
        industries   = [r[0] for r in conn.execute(
            "SELECT DISTINCT industry FROM companies WHERE industry IS NOT NULL AND industry!=''  ORDER BY industry LIMIT 60"
        ).fetchall()]
        signal_types = [r[0] for r in conn.execute(
            "SELECT DISTINCT signal_type FROM alerts_sent WHERE dry_run=0 ORDER BY signal_type"
        ).fetchall()]
        counts = {r[0]: r[1] for r in conn.execute(
            "SELECT signal_type, COUNT(*) cnt FROM alerts_sent WHERE dry_run=0 GROUP BY signal_type ORDER BY cnt DESC"
        ).fetchall()}
        total     = conn.execute("SELECT COUNT(*) FROM alerts_sent WHERE dry_run=0").fetchone()[0]
        companies = conn.execute("SELECT COUNT(DISTINCT apollo_id) FROM alerts_sent WHERE dry_run=0").fetchone()[0]
        conn.close()
        return jsonify({"industries": industries, "signal_types": signal_types,
                        "counts": counts, "total_signals": total, "total_companies": companies})
    except Exception as e:
        return jsonify({"error": str(e)})


@app.route("/api/insights/<account_id>")
@login_required
def insights_generate(account_id):
    import sqlite3, re as _re
    from pathlib import Path
    db_map = {"healthcare": Path(__file__).parent/"data"/"tracker.db",
              "csg":        Path(__file__).parent/"data"/"tracker_csg_v2.db"}
    db_path = db_map.get(account_id)
    if not db_path or not db_path.exists():
        return jsonify({"error": "Unknown account"}), 404
    api_key = os.environ.get("OPENAI_API_KEY","")
    if not api_key:
        return jsonify({"error": "OpenAI API key not configured"}), 500

    signal_types = request.args.getlist("signal_type")
    severities   = request.args.getlist("severity")
    days         = int(request.args.get("days", 90))
    industry     = request.args.get("industry","")

    try:
        conn = sqlite3.connect(str(db_path)); conn.row_factory = sqlite3.Row
        conds = ["a.dry_run = 0"]; params = []
        if signal_types:
            conds.append("a.signal_type IN (%s)" % ",".join("?"*len(signal_types))); params.extend(signal_types)
        if severities:
            conds.append("a.severity IN (%s)" % ",".join("?"*len(severities))); params.extend(severities)
        if days > 0:
            conds.append("a.signal_date >= date('now',?)"); params.append("-%d days" % days)
        if industry:
            conds.append("c.industry LIKE ?"); params.append("%"+industry+"%")
        where = " AND ".join(conds)
        rows = conn.execute(
            "SELECT c.name,c.domain,c.industry,c.city,c.state,"
            "a.signal_type,a.signal_detail,a.severity,a.signal_date "
            "FROM alerts_sent a JOIN companies c ON a.apollo_id=c.apollo_id "
            "WHERE "+where+" ORDER BY CASE a.severity WHEN 'HIGH' THEN 1 WHEN 'MEDIUM' THEN 2 ELSE 3 END,"
            "a.signal_date DESC LIMIT 200", params
        ).fetchall()
        conn.close()
        if not rows:
            return jsonify({"error": "No signals found for those filters."}), 200

        signals = [dict(r) for r in rows]
        n_sig = len(signals); n_co = len(set(s["name"] for s in signals))
        acct  = "Healthcare" if account_id == "healthcare" else "CSG"

        by_co = {}
        for s in signals:
            co = s["name"]
            if co not in by_co:
                by_co[co] = {"domain": s.get("domain",""), "industry": s.get("industry",""), "sigs":[]}
            by_co[co]["sigs"].append(s)

        from datetime import date as _date
        _today = _date.today()
        def _age_days(d):
            try:
                return (_today - _date.fromisoformat(str(d)[:10])).days
            except Exception:
                return 9999

        type_counts, ind_counts = {}, {}
        for s in signals:
            type_counts[s["signal_type"]] = type_counts.get(s["signal_type"], 0) + 1
            if s.get("industry"):
                ind_counts[s["industry"]] = ind_counts.get(s["industry"], 0) + 1

        ctx_lines = []
        multi_intent = 0
        for co, info in sorted(by_co.items(),
            key=lambda x: (-sum(1 for s in x[1]["sigs"] if s["severity"]=="HIGH"), -len(x[1]["sigs"])))[:80]:
            sigs = info["sigs"]
            stypes = sorted(set(s["signal_type"] for s in sigs))
            recent = sum(1 for s in sigs if _age_days(s["signal_date"]) <= 30)
            momentum = "RISING" if recent > len(sigs) - recent else ("ACTIVE" if recent else "COOLING")
            flags = []
            if len(stypes) >= 2:
                flags.append("MULTI-INTENT")
                multi_intent += 1
            if any(_age_days(s["signal_date"]) <= 7 for s in sigs):
                flags.append("FRESH<7d")
            sig_str = " | ".join(
                "%s(%s,%s)%s" % (s["signal_type"], s["severity"], s["signal_date"],
                    ": "+s["signal_detail"][:140] if s.get("signal_detail") else "")
                for s in sigs[:6])
            ctx_lines.append("[%s | %s | %s] %d sigs, %s%s — %s" % (
                co, info["domain"], info["industry"], len(sigs), momentum,
                (" " + ",".join(flags)) if flags else "", sig_str))

        stats_lines = (
            "DATASET STATS: %d signals across %d companies. "
            "Signal-type distribution: %s. Top industries: %s. "
            "Multi-intent companies (2+ distinct signal types): %d."
            % (n_sig, n_co,
               ", ".join("%s:%d" % kv for kv in sorted(type_counts.items(), key=lambda x: -x[1])),
               ", ".join("%s:%d" % kv for kv in sorted(ind_counts.items(), key=lambda x: -x[1])[:8]),
               multi_intent))

        schema = (
            '{"headline":"one punchy 8-12 word headline capturing this week in the market",'
            +'  "brief":"3-sentence leadership brief naming hottest companies, dominant signal pattern, ONE sales action.",'
            +'  "kairo_take":"one bold, non-obvious strategic observation from the data that a human analyst would likely miss",'
            +'  "week_priority":[{"rank":1,"company":"","domain":"","signal":"specific signal","pitch":"exact service+why","service":"SEO|PPC|Content|Brand|RevOps","call_timing":"Call today|Call this week|Warm email first","hook":"one-line conversation opener citing the signal"}],'
            +'  "market_pulse":["specific data-backed observation citing companies"],'
            +'  "strategic_moves":[{"move":"title","rationale":"signal-backed reason","impact":"qualitative business impact, no dollar figures","owner":"BDR|Account Exec|Marketing|Leadership"}],'
            +'  "pipeline":[{"name":"","domain":"","intent_score":85,"momentum":"rising|steady|cooling","signals":["type"],"why_now":"","service_fit":["SEO"],"contact_title":"best job title to approach","hook":"one-line opener citing their signal","next_step":""}],'
            +'  "actions":[{"rank":1,"type":"outreach","company":"","action":"","reason":"","deadline":"Today","urgency":"HIGH"}],'
            +'  "outreach":[{"company":"","domain":"","timing":"now","signal_hook":"","subject":"","opening":"","talking_points":[""],"cta":""}],'
            +'  "themes":[{"theme":"","count":0,"companies":[""],"campaign_angle":""}],'
            +'  "risks":[{"company":"","flag":"","implication":""}]}'
        )

        system_prompt = (
            "You are Kairo, Position2's elite revenue-intelligence AI. Position2 is a B2B digital "
            "marketing agency. Services: SEO & Organic Growth | Performance Marketing "
            "(Google/Meta/LinkedIn Ads) | Content Strategy | Brand & Website | Revenue Operations & HubSpot. "
            "You brief the CEO and Head of Sales on THIS WEEK's pipeline priorities. "
            "METHOD — reason through these steps before writing: "
            "(1) Weight every signal: severity (HIGH=3, MEDIUM=2, LOW=1) x recency (<7d x2, <30d x1.5, older x1); "
            "MULTI-INTENT companies (2+ distinct signal types) are the strongest buying-window evidence. "
            "(2) Score intent 0-100 from that weighting and be honest: most companies belong at 30-70; reserve 85+ "
            "for multi-intent + HIGH + fresh. Momentum flags in the data (RISING/ACTIVE/COOLING) must drive the "
            "pipeline momentum field. "
            "(3) Hunt cross-company patterns: sector waves, leadership migrations between tracked companies, funding "
            "clusters in one niche, timing coincidences. These power market_pulse, themes and kairo_take. "
            "(4) For each company, reason WHY the signal opens a marketing-services window NOW: new CMO/CEO = vendor "
            "review window (~90 days); funding = growth mandate and paid-media budget unlock; M&A = brand and website "
            "consolidation work; IPO = scrutiny on organic visibility and analyst-facing content; expansion/news = "
            "momentum to amplify. Map each to the single best-fit Position2 service. "
            "(5) Separate INTERNAL fields from PROSPECT-FACING copy. Internal fields (signal, why_now, reason, "
            "rationale, pitch, impact) may cite signal types and dates. PROSPECT-FACING copy (hook, subject, opening, "
            "talking_points, cta) is what a rep would actually say or send: NEVER mention dates, the word 'signal', "
            "or anything implying we monitor the company ('I saw', 'I noticed', 'your May 13 announcement'). Refer "
            "to public events naturally and obliquely ('as the new facility comes online'). Lead with their problem, "
            "sound human, zero buzzwords, no exclamation marks. "
            "BANNED: generic filler ('great fit', 'reach out to discuss', 'leverage synergies'), invented facts, and "
            "ANY revenue estimates, pipeline values, or dollar figures. Every claim must trace to a signal in the "
            "data. Specific beats clever; concise beats long. "
            "Return ONLY valid JSON exactly matching this schema: "+schema+" "
            "RULES: week_priority=top 6 by urgency; pipeline=top 14 scored 0-100 with honest momentum — include mid "
            "and lower-score watchlist companies too, not only the hot ones; actions=6 ranked; outreach=8 "
            "personalised with <55-char human, curiosity-driven subjects (no spammy caps); themes=4 each with a "
            "usable campaign_angle; risks=2-3 only if real. kairo_take must be a genuinely non-obvious pattern, "
            "never a summary."
        )

        from openai import OpenAI
        oai  = OpenAI(api_key=api_key, timeout=120.0, max_retries=1)
        user_msg = (
            "Analyse %d signals from %d %s-market companies.\n%s\n\n"
            "COMPANY SIGNAL DATA (format: [name | domain | industry] n sigs, MOMENTUM FLAGS — "
            "type(severity,date): detail):\n\n%s\n\nBrief the CEO. Respond with the JSON object only."
            % (n_sig, n_co, acct, stats_lines, "\n".join(ctx_lines)))
        raw, _used_model = _kairo_chat_json(oai, [
            {"role":"system","content":system_prompt},
            {"role":"user","content":user_msg}
        ], 6000)
        if "```" in raw:
            m = _re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
            raw = m.group(1).strip() if m else raw
        s2=raw.find("{"); e2=raw.rfind("}")
        if s2!=-1 and e2!=-1: raw=raw[s2:e2+1]
        # Handle truncated JSON by trying progressively shorter strings
        insights = None
        for attempt in [raw, raw[:raw.rfind("},")+1]+"}" if "}," in raw else raw]:
            try:
                insights = json.loads(attempt)
                break
            except json.JSONDecodeError:
                pass
        if insights is None:
            # Last resort: parse up to last valid closing brace
            for i in range(len(raw)-1, 0, -1):
                if raw[i]=="}":
                    try:
                        insights = json.loads(raw[:i+1])
                        break
                    except Exception:
                        continue
        if insights is None:
            return jsonify({"error": "GPT returned invalid JSON. Try again."})
        insights = _strip_revenue_fields(insights)
        return jsonify({"ok":True,"signals_analyzed":n_sig,"companies_analyzed":n_co,"model":_used_model,"insights":insights})
    except Exception as e:
        import traceback; log.error("insights_generate: %s", traceback.format_exc())
        return jsonify({"error": str(e)})


@app.route("/api/ppc-chat-debug")
@login_required
def ppc_chat_debug():
    """Shows exactly what data the chatbot sees — use to diagnose blank/wrong answers."""
    _PPC_CTX_CACHE["ts"] = 0   # force refresh
    ctx = _build_ppc_context()
    return f"<pre style='font-size:12px;padding:20px'>{ctx}</pre>", 200


# ── Ad Intelligence data helper (for chatbot) ────────────────────────────────
def get_ad_intelligence_data(competitor=None, ad_format=None, status=None,
                              keyword=None, limit=50):
    """
    Fetch competitor ad data from the Ad Intelligence Google Sheet via service account.
    The sheet must be shared with signal-tracker@signal-tracker-496308.iam.gserviceaccount.com

    Args:
        competitor : filter by competitor name or domain
        ad_format  : 'image', 'text', or 'video'
        status     : 'active' or 'inactive'
        keyword    : search in headline / description / keywords
        limit      : max rows to return (default 50)
    """
    try:
        svc = _sheets_service()
        # Read header row first to map columns robustly
        r_hdr = svc.spreadsheets().values().get(
            spreadsheetId=AD_INTEL_SHEET_ID, range="A1:AH1").execute()
        headers = [c.strip() for c in (r_hdr.get("values") or [[]])[0]]

        r_data = svc.spreadsheets().values().get(
            spreadsheetId=AD_INTEL_SHEET_ID, range="A2:AH2000").execute()
        data_rows = r_data.get("values") or []

    except Exception as e:
        err = str(e)
        if "403" in err or "permission" in err.lower() or "not found" in err.lower():
            return {
                "status": "error",
                "error": "Ad Intelligence sheet not shared with the service account.",
                "fix": "Share Google Sheet ID 16U5_QSxMmrAGKvK5dHScBu1Et4BJ1p8Q1ns5LycRA0s "
                       "with signal-tracker@signal-tracker-496308.iam.gserviceaccount.com (Viewer access).",
            }
        return {"status": "error", "error": err}

    def _v(row, col_name):
        try:
            idx = headers.index(col_name)
            return row[idx] if idx < len(row) else ""
        except ValueError:
            return ""

    ads = []
    for row in data_rows:
        domain = _v(row, "Domain")
        if not domain or domain == "Domain":
            continue
        ads.append({
            "competitor":      _v(row, "Advertiser Name") or domain,
            "domain":          domain,
            "format":          _v(row, "Format"),
            "platform":        _v(row, "Platform"),
            "status":          _v(row, "Status"),
            "headline":        _v(row, "Headline"),
            "description":     _v(row, "Description"),
            "full_text":       _v(row, "Full Ad Text"),
            "cta":             _v(row, "CTA"),
            "keywords":        _v(row, "Keywords"),
            "messaging_angle": _v(row, "Messaging Angle"),
            "value_prop":      _v(row, "Value Proposition"),
            "offer":           _v(row, "Offer"),
            "first_shown":     _v(row, "First Shown"),
            "last_shown":      _v(row, "Last Shown"),
            "regions":         _v(row, "Regions Served"),
        })

    total_before = len(ads)

    if competitor:
        c = competitor.lower()
        ads = [a for a in ads if c in a["domain"].lower() or c in a["competitor"].lower()]
    if ad_format:
        ads = [a for a in ads if a["format"].lower() == ad_format.lower()]
    if status:
        ads = [a for a in ads if a["status"].lower() == status.lower()]
    if keyword:
        kw = keyword.lower()
        ads = [a for a in ads if
               kw in a["headline"].lower() or kw in a["description"].lower()
               or kw in a["full_text"].lower() or kw in a["keywords"].lower()]

    format_counts  = dict(Counter(a["format"]  for a in ads if a["format"]).most_common())
    status_counts  = dict(Counter(a["status"]  for a in ads if a["status"]).most_common())
    comp_counts    = dict(Counter(a["competitor"] for a in ads if a["competitor"]).most_common())
    top_ctas       = [c for c, _ in Counter(a["cta"] for a in ads if a["cta"] and len(a["cta"]) < 50).most_common(5)]
    top_keywords   = [k.strip() for k, _ in Counter(
        kw.strip() for a in ads for kw in a["keywords"].split(",") if kw.strip()
    ).most_common(10)]

    return {
        "status": "ok",
        "total_in_sheet": total_before,
        "total_matching_filters": len(ads),
        "returned": min(len(ads), limit),
        "by_competitor": comp_counts,
        "by_format": format_counts,
        "by_status": status_counts,
        "top_ctas": top_ctas,
        "top_keywords": top_keywords,
        "ads": [
            {k: v for k, v in a.items() if k != "full_text"}
            for a in ads[:limit]
        ],
    }



@app.route("/api/company-analysis/<account_id>")
@login_required
def company_analysis(account_id):
    """Deep AI analysis of a single company for the insights drawer."""
    import sqlite3, re as _re
    from pathlib import Path
    db_map = {"healthcare": Path(__file__).parent/"data"/"tracker.db",
              "csg":        Path(__file__).parent/"data"/"tracker_csg_v2.db"}
    db_path = db_map.get(account_id)
    if not db_path or not db_path.exists():
        return jsonify({"error": "Unknown account"}), 404
    api_key = os.environ.get("OPENAI_API_KEY", "")
    if not api_key:
        return jsonify({"error": "OpenAI API key not configured"}), 500
    company_name = request.args.get("company", "")
    if not company_name:
        return jsonify({"error": "company parameter required"}), 400
    try:
        conn = sqlite3.connect(str(db_path)); conn.row_factory = sqlite3.Row
        rows = conn.execute(
            "SELECT a.signal_type, a.signal_detail, a.severity, a.signal_date, "
            "c.industry, c.domain, c.city, c.state "
            "FROM alerts_sent a JOIN companies c ON a.apollo_id=c.apollo_id "
            "WHERE c.name LIKE ? AND a.dry_run=0 ORDER BY a.signal_date DESC LIMIT 20",
            ["%" + company_name + "%"]
        ).fetchall()
        conn.close()
        if not rows:
            return jsonify({"error": "No signals found for this company"}), 200
        signals = [dict(r) for r in rows]
        sig_str = "\n".join(
            "- %s (%s, %s)%s" % (s["signal_type"], s["severity"], s["signal_date"],
                ": "+s["signal_detail"][:160] if s.get("signal_detail") else "")
            for s in signals)
        industry = signals[0].get("industry","") if signals else ""
        co_domain = signals[0].get("domain","") if signals else ""
        co_loc = ", ".join(x for x in [signals[0].get("city") or "", signals[0].get("state") or ""] if x) if signals else ""
        system = (
            "You are Kairo, senior B2B sales strategist at Position2 (SEO & Organic Growth, PPC/Performance "
            "Marketing, Content Strategy, Brand & Website, RevOps & HubSpot). Build a rigorous, signal-grounded "
            "prospect analysis. Reason first: what do the signals (their types, severity, dates, and sequence) "
            "imply about budget timing, internal change, and marketing gaps? Score honestly — most prospects are "
            "40-75; reserve 85+ for multiple fresh HIGH signals. talking_points, subject_lines and email_opening are "
            "PROSPECT-FACING: ground them in the signals but NEVER cite dates, the word 'signal', or anything that "
            "sounds like surveillance ('I saw', 'I noticed', 'your May 13 announcement') - refer to public events "
            "naturally and obliquely, lead with their problem, zero buzzwords, no exclamation marks. score_reason, "
            "why_now and urgency_reason are INTERNAL: dates allowed there. Objections must be the realistic ones "
            "for this industry. Subject lines: human, specific, curiosity-driven, <55 chars, no clickbait caps. "
            "NEVER include revenue estimates or dollar figures. "
            "Return ONLY valid JSON: "
            '{"score":85,"score_reason":"one sentence why",'
            '"company_context":"2 sentences about what this company does and why they matter",'
            '"why_now":"2 sentences on why right now is the perfect time to reach out",'
            '"talking_points":["specific point 1 tied to signal","specific point 2","specific point 3"],'
            '"objections":[{"objection":"likely pushback","response":"how to handle it"}],'
            '"subject_lines":["option 1 <55 chars","option 2","option 3"],'
            '"email_opening":"2 sentence personalized opening referencing their specific situation",'
            '"recommended_service":"SEO|PPC|Content|Brand|RevOps",'
            '"service_reason":"why this specific service fits their situation",'
            '"urgency":"HIGH|MEDIUM|LOW","urgency_reason":"why"}'
        )
        from openai import OpenAI
        oai = OpenAI(api_key=api_key, timeout=90.0, max_retries=1)
        raw, _used_model = _kairo_chat_json(oai, [
            {"role": "system", "content": system},
            {"role": "user", "content": "Company: %s\nDomain: %s\nIndustry: %s\nLocation: %s\nSignals (newest first):\n%s"
                % (company_name, co_domain, industry, co_loc, sig_str)}
        ], 1600)
        if "```" in raw:
            m = _re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
            raw = m.group(1).strip() if m else raw
        s2=raw.find("{"); e2=raw.rfind("}")
        if s2!=-1 and e2!=-1: raw=raw[s2:e2+1]
        return jsonify({"ok": True, "company": company_name, "model": _used_model,
                        "analysis": _strip_revenue_fields(json.loads(raw))})
    except Exception as e:
        return jsonify({"error": str(e)})



@app.route("/api/generate-email/<account_id>")
@login_required
def generate_email(account_id):
    """GPT-powered personalised email using company signals."""
    import sqlite3, re as _re
    from pathlib import Path
    db_map = {"healthcare": Path(__file__).parent/"data"/"tracker.db",
              "csg":        Path(__file__).parent/"data"/"tracker_csg_v2.db"}
    db_path = db_map.get(account_id)
    if not db_path or not db_path.exists():
        return jsonify({"error": "Unknown account"})
    api_key = os.environ.get("OPENAI_API_KEY","")
    if not api_key:
        return jsonify({"error": "OpenAI API key not configured"})
    company = request.args.get("company","").strip()
    service = request.args.get("service","")
    tone    = (request.args.get("tone","") or "direct").strip().lower()
    if not company:
        return jsonify({"error": "company parameter required"})
    try:
        conn = sqlite3.connect(str(db_path)); conn.row_factory = sqlite3.Row
        rows = conn.execute(
            "SELECT a.signal_type,a.signal_detail,a.severity,a.signal_date,"
            "c.industry,c.domain FROM alerts_sent a "
            "JOIN companies c ON a.apollo_id=c.apollo_id "
            "WHERE c.name LIKE ? AND a.dry_run=0 ORDER BY a.signal_date DESC LIMIT 15",
            ["%"+company+"%"]
        ).fetchall()
        conn.close()
        signals = [dict(r) for r in rows]
        if not signals:
            return jsonify({"error": "No signals found for " + company})

        industry = signals[0].get("industry","") if signals else ""
        domain   = signals[0].get("domain","")   if signals else ""
        sig_lines = "\n".join(
            "- %s (%s) on %s%s" % (
                s["signal_type"], s["severity"], s["signal_date"],
                ": "+s["signal_detail"][:120] if s.get("signal_detail") else ""
            ) for s in signals[:10]
        )

        tone_guide = {
            "direct":    "TONE: confident and direct, zero fluff - a sharp consultant who respects the reader's time.",
            "warm":      "TONE: warm and human, lightly conversational, still professional.",
            "executive": "TONE: senior executive to senior executive - measured, strategic, no casual phrases.",
        }.get(tone, "TONE: confident and direct, zero fluff.")

        system = (
            "You are Kairo, writing outreach for Position2, a B2B digital marketing agency "
            "(SEO | Performance Marketing/PPC | Content Strategy | Brand & Website | Revenue Operations). "
            "Write an email a thoughtful senior consultant would actually send - never anything that smells of "
            "AI or mail-merge.\n\n"
            "The signals provided are INTERNAL intelligence. Use them ONLY to understand the company situation.\n"
            "HARD RULES:\n"
            "- NEVER mention dates, the word signal, announcements you saw or noticed, or anything implying we "
            "monitor them. Banned openers: I saw / I noticed / I came across / Congratulations on / Hope this "
            "finds you well / Quick question.\n"
            "- Refer to public events only obliquely and naturally (as the new facility comes online; with the "
            "team growing) - no dates, no press-release specifics.\n"
            "- The first sentence must be about THEIR world - a real problem or opportunity - and it must earn "
            "the second sentence. Never open with us.\n"
            "- Include one concrete, useful observation or idea they could act on even without replying. That "
            "is the value of the email.\n"
            "- Mention Position2 once, briefly, as credibility - no service list, no we-help-companies-like-you.\n"
            "- ONE low-friction CTA phrased as an easy yes/no question. Never hop-on-a-call-to-discuss.\n"
            "- Under 110 words across opening+body+cta. Short sentences. 7th-grade readability. No buzzwords "
            "(leverage, synergies, streamline, elevate, unlock, empower, seamless, cutting-edge) and no "
            "exclamation marks.\n"
            "- subject: under 45 characters, natural and specific, sentence case, no clickbait.\n"
            "- greeting: exactly Hi {FirstName}, so the rep can personalise.\n"
            "- ps: only if there is a genuinely useful extra thought, otherwise an empty string. Never a second "
            "pitch.\n"
            "- Never invent facts, metrics, client names, or revenue/dollar figures.\n"
            + tone_guide + "\n\n"
            "Return ONLY valid JSON:\n"
            '{"subject":"","greeting":"Hi {FirstName},","opening":"<1 sentence about their world>",'
            '"body":"<2-3 short sentences: useful insight, then one line of Position2 credibility>",'
            '"cta":"<one easy yes/no question>","ps":"",'
            '"recommended_service":"<SEO|PPC|Content|Brand|RevOps>",'
            '"why_now":"<INTERNAL rep note on timing - dates allowed here; one sentence>"}'
        )

        user_msg = "Company: %s\nIndustry: %s\nDomain: %s\nPreferred service: %s\n\nSignals:\n%s\n\nWrite the email." % (
            company, industry, domain, service or "best fit", sig_lines)

        from openai import OpenAI
        oai = OpenAI(api_key=api_key)
        raw, _m = _kairo_completion(
            oai, [{"role":"system","content":system},{"role":"user","content":user_msg}], 800)
        if "```" in raw:
            m = _re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
            raw = m.group(1).strip() if m else raw
        s2=raw.find("{"); e2=raw.rfind("}")
        if s2!=-1 and e2!=-1: raw=raw[s2:e2+1]
        email_data = json.loads(raw)
        return jsonify({"ok":True,"company":company,"email":email_data,"signals_used":len(signals)})
    except Exception as e:
        import traceback; log.error("generate_email: %s", traceback.format_exc())
        return jsonify({"error": str(e)})


@app.route("/api/research-company/<account_id>")
@login_required
def research_company(account_id):
    """AI research on a company: GPT + web search -> key facts, insights, Position2 angle."""
    import sqlite3, re as _re
    from pathlib import Path
    db_map = {"healthcare": Path(__file__).parent/"data"/"tracker.db",
              "csg":        Path(__file__).parent/"data"/"tracker_csg_v2.db"}
    db_path = db_map.get(account_id)
    if not db_path or not db_path.exists():
        return jsonify({"error": "Unknown account"})
    api_key = os.environ.get("OPENAI_API_KEY","")
    if not api_key:
        return jsonify({"error": "OpenAI API key not configured"})
    company = request.args.get("company","").strip()
    domain  = request.args.get("domain","").strip()
    if not company:
        return jsonify({"error": "company parameter required"})
    try:
        # Pull known signals for grounding context
        sig_lines = ""
        try:
            conn = sqlite3.connect(str(db_path)); conn.row_factory = sqlite3.Row
            rows = conn.execute(
                "SELECT a.signal_type,a.signal_detail,a.severity,a.signal_date,c.industry,c.domain "
                "FROM alerts_sent a JOIN companies c ON a.apollo_id=c.apollo_id "
                "WHERE c.name LIKE ? AND a.dry_run=0 ORDER BY a.signal_date DESC LIMIT 10",
                ["%"+company+"%"]).fetchall()
            conn.close()
            sigs = [dict(r) for r in rows]
            if sigs and not domain:
                domain = sigs[0].get("domain","") or ""
            sig_lines = "\n".join(
                "- %s (%s) on %s%s" % (s["signal_type"], s["severity"], s["signal_date"],
                    ": "+s["signal_detail"][:120] if s.get("signal_detail") else "")
                for s in sigs)
        except Exception:
            pass

        system = (
            "You are Kairo, Position2’s sales-intelligence research AI. Position2 is a digital marketing agency. "
            "Services: SEO & Organic Growth | Performance Marketing (Google/Meta/LinkedIn Ads) | "
            "Content Strategy | Brand & Website | Revenue Operations & HubSpot. "
            "Research the given company using web search. Find what they do, recent news, "
            "leadership, market position, and their likely digital-marketing gaps. "
            "NEVER include revenue estimates or dollar figures. "
            "Return ONLY valid JSON (no markdown):\n"
            '{"overview":"2-3 sentences on what the company does and their market position",'
            '"recent_developments":[{"date":"YYYY-MM or recent","headline":"","detail":"1 sentence"}],'
            '"key_people":[{"name":"","role":""}],'
            '"digital_presence":"1-2 sentences on their website/SEO/ads/social footprint and visible gaps",'
            '"opportunities":["specific marketing gap or opportunity Position2 could address"],'
            '"position2_angle":"2 sentences: which Position2 services fit and why, tied to findings",'
            '"recommended_services":["SEO|PPC|Content|Brand|RevOps"],'
            '"conversation_starters":["natural human opener grounded in a real finding - no dates, never sounding like surveillance"],'
            '"sources":[{"title":"","url":""}]}'
        )
        user_msg = "Research this company NOW:\nCompany: %s\nDomain: %s\n%s" % (
            company, domain or "unknown",
            "Known signals from our tracker:\n"+sig_lines if sig_lines else "")

        from openai import OpenAI
        oai   = OpenAI(api_key=api_key)
        _msgs = [{"role":"system","content":system},{"role":"user","content":user_msg}]
        model = os.environ.get("OPENAI_INSIGHTS_MODEL") or "gpt-5.4"
        raw, web_used = _responses_web_search(oai, model, _msgs, 2500)
        if not raw:
            model = os.environ.get("OPENAI_MODEL","gpt-4o-mini")
            raw, web_used = _responses_web_search(oai, model, _msgs, 2500)
        # Fallback: plain completion using only tracker signals
        if not raw:
            resp = oai.chat.completions.create(
                model=model,
                messages=[{"role":"system","content":system.replace("using web search","using your knowledge (web search unavailable)")},
                          {"role":"user","content":user_msg}],
                max_completion_tokens=2000,
            )
            raw = resp.choices[0].message.content.strip()
        if "```" in raw:
            m = _re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
            raw = m.group(1).strip() if m else raw
        s2=raw.find("{"); e2=raw.rfind("}")
        if s2!=-1 and e2!=-1: raw=raw[s2:e2+1]
        research = _strip_revenue_fields(json.loads(raw))
        return jsonify({"ok": True, "company": company, "domain": domain,
                        "web_search_used": web_used, "research": research})
    except Exception as e:
        import traceback; log.error("research_company: %s", traceback.format_exc())
        return jsonify({"error": str(e)})


@app.route("/api/kairo-chat/<account_id>", methods=["POST"])
@login_required
def kairo_chat(account_id):
    """Conversational Kairo: grounded on the account signal DB, web-search for the rest."""
    import sqlite3, re as _re
    from pathlib import Path
    db_map = {"healthcare": Path(__file__).parent/"data"/"tracker.db",
              "csg":        Path(__file__).parent/"data"/"tracker_csg_v2.db"}
    db_path = db_map.get(account_id)
    if not db_path or not db_path.exists():
        return jsonify({"error": "Unknown account"})
    api_key = os.environ.get("OPENAI_API_KEY", "")
    if not api_key:
        return jsonify({"error": "OpenAI API key not configured"})
    body = request.get_json(silent=True) or {}
    question = (body.get("question") or "").strip()
    history = body.get("history") or []
    files = body.get("files") or []
    if not question:
        return jsonify({"error": "empty question"})

    img_files = [f for f in files if isinstance(f, dict) and f.get("image") and f.get("base64")][:4]
    txt_attached = [f for f in files if isinstance(f, dict) and not f.get("image")]

    # Attached-file context (extracted client-side via /api/ppc-upload)
    att_ctx = ""
    if txt_attached:
        chunks = []
        for fobj in txt_attached[:4]:
            nm = str(fobj.get("name") or "file")[:120]
            ct = str(fobj.get("content") or "")[:12000]
            if ct.strip():
                chunks.append("=== ATTACHED FILE: %s ===\n%s" % (nm, ct))
        if chunks:
            att_ctx = ("\n\nATTACHED FILES (uploaded by the user — treat as primary context; quote and "
                       "analyse their actual contents):\n" + "\n\n".join(chunks))

    # Detect a requested output format (csv / xlsx / docx / pdf / pptx)
    export_format = ""
    _ql = question.lower()
    if _re.search(r"\b(as|in|into|to|export|download|give|create|make|generate|build|produce|format|convert)\b", _ql):
        for _f, _pat in (("csv", r"\bcsv\b"), ("xlsx", r"\b(xlsx|xls|excel|spreadsheet)\b"),
                         ("docx", r"\b(docx|word)\b"), ("pdf", r"\bpdf\b"),
                         ("pptx", r"\b(pptx|ppt|powerpoint|slide deck|slides|deck)\b")):
            if _re.search(_pat, _ql):
                export_format = _f
                break
    try:
        conn = sqlite3.connect(str(db_path)); conn.row_factory = sqlite3.Row
        counts = {r[0]: r[1] for r in conn.execute(
            "SELECT signal_type, COUNT(*) FROM alerts_sent WHERE dry_run=0 GROUP BY signal_type")}
        total_sig = sum(counts.values())
        total_co = conn.execute("SELECT COUNT(DISTINCT apollo_id) FROM alerts_sent WHERE dry_run=0").fetchone()[0]
        acct_label = "Healthcare" if account_id == "healthcare" else "CSG"
        ql = question.lower()
        names = [r[0] for r in conn.execute("SELECT DISTINCT name FROM companies WHERE name IS NOT NULL")]
        matched = [n for n in names if n and len(n) > 2 and n.lower() in ql][:6]
        ctx = []
        if matched:
            for nm in matched:
                rows = conn.execute(
                    "SELECT a.signal_type,a.severity,a.signal_date,a.signal_detail,c.domain,c.industry "
                    "FROM alerts_sent a JOIN companies c ON a.apollo_id=c.apollo_id "
                    "WHERE c.name=? AND a.dry_run=0 ORDER BY a.signal_date DESC LIMIT 12", [nm]).fetchall()
                if rows:
                    sl = " | ".join("%s(%s,%s)%s" % (
                        r["signal_type"], r["severity"], r["signal_date"],
                        (": " + r["signal_detail"][:80]) if r["signal_detail"] else "") for r in rows[:8])
                    ctx.append("[%s | %s | %s] %s" % (nm, rows[0]["domain"], rows[0]["industry"], sl))
        else:
            rows = conn.execute(
                "SELECT c.name,c.domain,COUNT(*) n,SUM(CASE WHEN a.severity='HIGH' THEN 1 ELSE 0 END) hi "
                "FROM alerts_sent a JOIN companies c ON a.apollo_id=c.apollo_id "
                "WHERE a.dry_run=0 GROUP BY c.apollo_id ORDER BY hi DESC,n DESC LIMIT 12").fetchall()
            for r in rows:
                ctx.append("%s (%s) - %d signals, %d HIGH" % (r["name"], r["domain"], r["n"], r["hi"]))
        conn.close()

        overview = "Account: %s market. %d tracked signals across %d companies. Signal mix: %s." % (
            acct_label, total_sig, total_co,
            ", ".join("%s %d" % (k, v) for k, v in sorted(counts.items(), key=lambda x: -x[1])))
        ctx_str = "\n".join(ctx) or "(no specific company matched - use the overview and web search)"
        system = (
            "You are Kairo, Position2's signal-intelligence assistant. Position2 is a B2B digital marketing "
            "agency (SEO, PPC, Content, Brand & Website, RevOps). Answer the user accurately and concisely. "
            "Use the ACCOUNT SIGNAL DATA below for questions about tracked companies and signals; use web search "
            "for company research, recent news, people, contacts, or anything not in the data. If asked to draft an "
            "email or message, make it tight, personalised and HUMAN: never cite signal dates or imply we monitor "
            "the company ('I saw your May 13 announcement') in prospect-facing copy - refer to public events "
            "naturally and obliquely. Never invent revenue or dollar figures. Cite specific "
            "companies and signals. Format with short, clean markdown (bold, links, short lists). "
            "If files, images or screenshots are attached, ground your answer in their ACTUAL contents — quote real numbers, names and rows "
            "from them, and combine them with signal data where relevant. "
            "If the user asks for output as a file or specific format (CSV, Excel/XLSX, Word/DOCX, PDF, "
            "PowerPoint/PPTX, or a table), produce the COMPLETE content in clean markdown: a proper markdown table "
            "for tabular data, headings (#, ##) to structure documents and slides. The platform converts your "
            "markdown into the requested file, so never refuse a format and never truncate with placeholders.\n\n"
            "ACCOUNT OVERVIEW: %s\n\nRELEVANT SIGNAL DATA:\n%s%s" % (overview, ctx_str, att_ctx))

        msgs = [{"role": "system", "content": system}]
        for m in history[-8:]:
            role = m.get("role")
            if role in ("user", "assistant") and m.get("content"):
                msgs.append({"role": role, "content": str(m["content"])[:2000]})
        msgs.append({"role": "user", "content": question})

        from openai import OpenAI
        oai = OpenAI(api_key=api_key, timeout=80.0, max_retries=1)
        _max_out = 2600 if (files or export_format) else 1100
        if img_files:
            # Vision path: send image(s) directly to the model as data URIs
            parts = [{"type": "text", "text": question}]
            for im in img_files:
                mime = str(im.get("mime") or "image/png")[:50]
                b64 = str(im.get("base64") or "")
                if not b64 or len(b64) > 9_000_000:
                    continue
                parts.append({"type": "image_url",
                              "image_url": {"url": "data:%s;base64,%s" % (mime, b64)}})
            msgs[-1] = {"role": "user", "content": parts}
            answer, _m = _kairo_completion(oai, msgs, _max_out)
            web = False
        else:
            model = os.environ.get("OPENAI_INSIGHTS_MODEL") or "gpt-5.4"
            answer, web = _responses_web_search(oai, model, msgs, _max_out)
            if not answer:
                answer, web = _responses_web_search(oai, os.environ.get("OPENAI_MODEL", "gpt-4o-mini"), msgs, _max_out)
            if not answer:
                answer, _m = _kairo_completion(oai, msgs, _max_out)
        return jsonify({"ok": True, "answer": answer, "web_search_used": web,
                        "export_format": export_format})
    except Exception as e:
        import traceback; log.error("kairo_chat: %s", traceback.format_exc())
        return jsonify({"error": str(e)})


# ── Kairo export: convert a markdown answer into a downloadable file ─────────

def _md_blocks(content):
    """Parse markdown-lite into (kind, payload) blocks: h1/h2/h3/li/p (str) and tr (list of cells)."""
    blocks = []
    for ln in content.splitlines():
        st = ln.strip()
        if not st:
            continue
        if st.startswith("|") and st.endswith("|") and st.count("|") >= 2:
            cells = [c.strip() for c in st.strip("|").split("|")]
            if all(set(c) <= set("-: ") for c in cells):
                continue  # separator row
            blocks.append(("tr", cells)); continue
        if st.startswith("### "): blocks.append(("h3", st[4:])); continue
        if st.startswith("## "):  blocks.append(("h2", st[3:])); continue
        if st.startswith("# "):   blocks.append(("h1", st[2:])); continue
        if st[:2] in ("- ", "* "): blocks.append(("li", st[2:])); continue
        blocks.append(("p", st))
    return blocks


def _md_table_rows(content):
    """First choice: markdown table rows. Fallback: CSV inside a code block."""
    rows = [v for k, v in _md_blocks(content) if k == "tr"]
    if rows:
        return rows
    import csv as _csv, io as _io, re as _re2
    m = _re2.search(r"```(?:csv)?\s*([\s\S]*?)```", content)
    blob = m.group(1).strip() if m else ""
    if blob and ("," in blob or "\t" in blob):
        return [r for r in _csv.reader(_io.StringIO(blob)) if any(x.strip() for x in r)]
    return []


def _strip_md(t):
    import re as _re2
    t = str(t)
    t = _re2.sub(r"\*\*([^*]+)\*\*", r"\1", t)
    t = _re2.sub(r"\[([^\]]+)\]\((https?://[^\s)]+)\)", r"\1 (\2)", t)
    return t.strip()


@app.route("/api/kairo-export", methods=["POST"])
@login_required
def kairo_export():
    """Convert Kairo markdown output into CSV / XLSX / DOCX / PDF / PPTX and stream it back."""
    import io
    from flask import send_file
    body = request.get_json(silent=True) or {}
    fmt = (body.get("format") or "").lower().strip()
    content = str(body.get("content") or "").strip()
    title = (body.get("title") or "Kairo Insights").strip()[:80] or "Kairo Insights"
    if not content:
        return jsonify({"error": "no content"}), 400
    if fmt not in ("csv", "xlsx", "docx", "pdf", "pptx"):
        return jsonify({"error": "unsupported format"}), 400
    fname = "kairo-insights." + fmt
    try:
        if fmt == "csv":
            import csv as _csv
            buf = io.StringIO()
            w = _csv.writer(buf)
            rows = _md_table_rows(content)
            if rows:
                for r in rows:
                    w.writerow([_strip_md(c) for c in r])
            else:
                for kind, val in _md_blocks(content):
                    w.writerow([_strip_md(val if isinstance(val, str) else " | ".join(val))])
            data = io.BytesIO(buf.getvalue().encode("utf-8-sig"))
            return send_file(data, mimetype="text/csv", as_attachment=True, download_name=fname)

        if fmt == "xlsx":
            import openpyxl
            from openpyxl.styles import Font
            wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Kairo"
            rows = _md_table_rows(content)
            if rows:
                for ri, r in enumerate(rows, 1):
                    for ci, c in enumerate(r, 1):
                        ws.cell(row=ri, column=ci, value=_strip_md(c))
                for c in ws[1]:
                    c.font = Font(bold=True)
            else:
                ri = 1
                for kind, val in _md_blocks(content):
                    cell = ws.cell(row=ri, column=1,
                                   value=_strip_md(val if isinstance(val, str) else " | ".join(val)))
                    if kind in ("h1", "h2", "h3"):
                        cell.font = Font(bold=True, size=13 if kind == "h1" else 12)
                    ri += 1
            for col in ws.columns:
                mx = max((len(str(c.value or "")) for c in col), default=10)
                ws.column_dimensions[col[0].column_letter].width = min(60, max(12, mx + 2))
            data = io.BytesIO(); wb.save(data); data.seek(0)
            return send_file(data, as_attachment=True, download_name=fname,
                mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

        if fmt == "docx":
            import docx
            doc = docx.Document()
            doc.add_heading(title, level=0)
            tbl = []
            def _flush():
                if not tbl:
                    return
                t = doc.add_table(rows=0, cols=max(len(r) for r in tbl))
                try: t.style = "Light Grid Accent 1"
                except Exception: pass
                for r in tbl:
                    cells = t.add_row().cells
                    for i2, c in enumerate(r):
                        if i2 < len(cells):
                            cells[i2].text = _strip_md(c)
                del tbl[:]
            for kind, val in _md_blocks(content):
                if kind == "tr":
                    tbl.append(val); continue
                _flush()
                if kind == "h1": doc.add_heading(_strip_md(val), level=1)
                elif kind == "h2": doc.add_heading(_strip_md(val), level=2)
                elif kind == "h3": doc.add_heading(_strip_md(val), level=3)
                elif kind == "li": doc.add_paragraph(_strip_md(val), style="List Bullet")
                else: doc.add_paragraph(_strip_md(val))
            _flush()
            data = io.BytesIO(); doc.save(data); data.seek(0)
            return send_file(data, as_attachment=True, download_name=fname,
                mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

        if fmt == "pptx":
            from pptx import Presentation
            prs = Presentation()
            slides, cur = [], [title, []]
            for kind, val in _md_blocks(content):
                txt = _strip_md(val if isinstance(val, str) else " | ".join(val))
                if kind in ("h1", "h2"):
                    if cur[1]:
                        slides.append(cur)
                    cur = [txt, []]
                else:
                    cur[1].append(("• " if kind == "li" else "") + txt)
            if cur[1] or not slides:
                slides.append(cur)
            for stitle, lines in slides[:30]:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = stitle[:90]
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                for i2, ln in enumerate(lines[:12]):
                    p = tf.paragraphs[0] if i2 == 0 else tf.add_paragraph()
                    p.text = ln[:180]
            data = io.BytesIO(); prs.save(data); data.seek(0)
            return send_file(data, as_attachment=True, download_name=fname,
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        # ── pdf ──
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        data = io.BytesIO()
        docp = SimpleDocTemplate(data, pagesize=A4, topMargin=18*mm, bottomMargin=18*mm)
        styles = getSampleStyleSheet()
        story = [Paragraph(title, styles["Title"]), Spacer(1, 6)]
        tbl = []
        def _flush_pdf():
            if not tbl:
                return
            t = Table([[_strip_md(c)[:90] for c in r] for r in tbl], hAlign="LEFT")
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4f46e5")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#cbd5e1")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f1f5f9")]),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(t); story.append(Spacer(1, 8)); del tbl[:]
        for kind, val in _md_blocks(content):
            if kind == "tr":
                tbl.append(val); continue
            _flush_pdf()
            txt = _strip_md(val).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            if kind == "h1": story.append(Paragraph(txt, styles["Heading1"]))
            elif kind == "h2": story.append(Paragraph(txt, styles["Heading2"]))
            elif kind == "h3": story.append(Paragraph(txt, styles["Heading3"]))
            elif kind == "li": story.append(Paragraph("• " + txt, styles["Normal"]))
            else: story.append(Paragraph(txt, styles["Normal"]))
        _flush_pdf()
        docp.build(story)
        data.seek(0)
        return send_file(data, mimetype="application/pdf", as_attachment=True, download_name=fname)
    except Exception as e:
        import traceback; log.error("kairo_export: %s", traceback.format_exc())
        return jsonify({"error": str(e)}), 500


@app.route("/api/refresh-dashboard", methods=["POST"])
@login_required
def refresh_dashboard():
    """Trigger the GitHub Action that fetches the latest signals (HIGH from
    Sheets, LOW from Google News with filters) for both accounts, rebuilds the
    dashboards (preserving Kairo), prunes news, and publishes."""
    token    = os.environ.get("GH_DISPATCH_TOKEN", "")
    repo     = os.environ.get("GH_REPO", "ai-positon2/intelligence-platform")
    workflow = os.environ.get("GH_WORKFLOW", "refresh-dashboards.yml")
    if not token:
        return jsonify({"error": "Refresh isn't wired up yet — add a GH_DISPATCH_TOKEN "
                                 "environment variable in Railway (a GitHub token with the "
                                 "'workflow' scope). Until then, use the manual commands below."}), 200
    try:
        url = "https://api.github.com/repos/%s/actions/workflows/%s/dispatches" % (repo, workflow)
        r = requests.post(url, json={"ref": "main"}, timeout=20, headers={
            "Authorization": "Bearer " + token,
            "Accept": "application/vnd.github+json",
            "X-GitHub-Api-Version": "2022-11-28",
        })
        if r.status_code in (201, 204):
            return jsonify({"ok": True,
                "message": "Refresh started. Kairo is fetching the latest HIGH signals (Sheets) and "
                           "LOW signals (Google News, filtered) for both accounts, rebuilding, and "
                           "publishing. Your dashboard updates automatically in a few minutes — reload then.",
                "actions_url": "https://github.com/%s/actions/workflows/%s" % (repo, workflow)})
        return jsonify({"error": "GitHub returned %d. Check the GH_DISPATCH_TOKEN scope/repo. %s"
                                 % (r.status_code, (r.text or "")[:160])}), 200
    except Exception as e:
        import traceback; log.error("refresh_dashboard: %s", traceback.format_exc())
        return jsonify({"error": str(e)}), 200


def _refresh_stage(name):
    n = (name or "").lower()
    if "fetch healthcare" in n: return "Fetching Healthcare signals (Sheets + News)…"
    if "fetch csg" in n:        return "Fetching CSG news…"
    if "rebuild" in n:          return "Rebuilding & scoring both accounts…"
    if "commit" in n or "publish" in n: return "Publishing…"
    if "restore" in n:          return "Rebuilding…"
    return "Preparing…"

@app.route("/api/refresh-status")
@login_required
def refresh_status():
    """Live status of the most recent refresh Action run (for the progress bar)."""
    import datetime as _dt, time as _t
    token    = os.environ.get("GH_DISPATCH_TOKEN", "")
    repo     = os.environ.get("GH_REPO", "ai-positon2/intelligence-platform")
    workflow = os.environ.get("GH_WORKFLOW", "refresh-dashboards.yml")
    if not token:
        return jsonify({"error": "Refresh isn't configured (missing GH_DISPATCH_TOKEN)."})
    def _epoch(iso):
        try:
            return _dt.datetime.strptime(iso, "%Y-%m-%dT%H:%M:%SZ").replace(tzinfo=_dt.timezone.utc).timestamp()
        except Exception:
            return 0.0
    try:
        since = float(request.args.get("since", 0) or 0)
        hdr = {"Authorization": "Bearer " + token, "Accept": "application/vnd.github+json",
               "X-GitHub-Api-Version": "2022-11-28"}
        rr = requests.get("https://api.github.com/repos/%s/actions/workflows/%s/runs?per_page=5" % (repo, workflow),
                          headers=hdr, timeout=15)
        if rr.status_code != 200:
            return jsonify({"error": "GitHub runs %d" % rr.status_code})
        run = None
        for w in rr.json().get("workflow_runs", []):
            if since <= 0 or _epoch(w.get("created_at", "")) >= since - 90:
                run = w; break
        if not run:
            return jsonify({"pending": True})
        status = run.get("status"); concl = run.get("conclusion")
        started = _epoch(run.get("run_started_at") or run.get("created_at") or "")
        elapsed = max(0, int(_t.time() - started)) if started else 0
        percent, stage = 8, "Queued…"
        if status == "completed":
            percent, stage = 100, ("Done" if concl == "success" else "Failed")
        else:
            steps = []
            jurl = run.get("jobs_url")
            if jurl:
                jr = requests.get(jurl, headers=hdr, timeout=15)
                if jr.status_code == 200:
                    jobs = jr.json().get("jobs", [])
                    if jobs: steps = jobs[0].get("steps", []) or []
            if steps:
                tot = len(steps); done = sum(1 for s in steps if s.get("status") == "completed")
                cur = next((s for s in steps if s.get("status") == "in_progress"), None)
                percent = min(95, int(round(100.0 * done / max(tot, 1))))
                stage = _refresh_stage(cur.get("name")) if cur else "Working…"
        return jsonify({"status": status, "conclusion": concl, "percent": percent,
                        "stage": stage, "elapsed": elapsed, "html_url": run.get("html_url", "")})
    except Exception as e:
        return jsonify({"error": str(e)})


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port)

# redeploy nudge 1781202493 (particles v2 + login music)

# redeploy nudge 1781210922 (remove orb + motif on hub/ppc/seo)

# redeploy nudge 20260613-162500
# redeploy nudge 20260612-075817
